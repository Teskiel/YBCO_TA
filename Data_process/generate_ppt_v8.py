# -*- coding: utf-8 -*-
"""
YBCO KID 五谐振器对比表征简报 v8 (final)。

数据源：
- R1-R5 各含完整 scraps cmplxIQ 分析（f0, Qi, S21 overlay, res shift, responsivity）
- compare_v2/ 五谐振器对比汇总图
- 01_resonance_detection/ 频谱总览
"""

from pathlib import Path
from PIL import Image as PILImage
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

SCRIPT_DIR = Path(__file__).resolve().parent
MERGED = SCRIPT_DIR / "output" / "merged"
CMP = MERGED / "compare_v2"
DETECT = MERGED / "01_resonance_detection"

RESONATORS = [
    {"name": "R1", "freq": "3.846 GHz", "freq_short": "3.846", "dip": "-5.80"},
    {"name": "R2", "freq": "4.010 GHz", "freq_short": "4.010", "dip": "-4.33"},
    {"name": "R3", "freq": "4.500 GHz", "freq_short": "4.500", "dip": "-5.98"},
    {"name": "R4", "freq": "4.997 GHz", "freq_short": "4.997", "dip": "-17.74"},
    {"name": "R5", "freq": "5.252 GHz", "freq_short": "5.252", "dip": "-6.10"},
]
R_DIRS = [MERGED / f"R{i+1}_{r['freq_short']}GHz" for i, r in enumerate(RESONATORS)]

# 为每个谐振器预设 4 个代表温度点 (使用 actual temp 名)
# 低温: ~6K, 中低: ~20K, 中高: ~40K, 高温: ~77K
TEMP_POINTS = ["5.991K", "19.977K", "39.825K", "76.204K"]
TEMP_LABELS = ["6 K", "20 K", "40 K", "77 K"]

PPTX_OUT = SCRIPT_DIR / "output" / "YBCO_KID_五谐振器对比_v8_final.pptx"

WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x1A, 0x1A, 0x1A)
GRAY   = RGBColor(0x66, 0x66, 0x66)
LGRAY  = RGBColor(0xAA, 0xAA, 0xAA)
ACCENT = RGBColor(0x1F, 0x77, 0xB4)
ACCENT2 = RGBColor(0xD6, 0x27, 0x28)
RCOLORS = [
    RGBColor(0x1F, 0x77, 0xB4),
    RGBColor(0xD6, 0x27, 0x28),
    RGBColor(0x2C, 0xA0, 0x2C),
    RGBColor(0xFF, 0x7F, 0x0E),
    RGBColor(0x94, 0x67, 0xBD),
]

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

_image_cache = {}
def _ar(p):
    p = Path(p)
    if not p.exists(): return None
    k = str(p)
    if k not in _image_cache:
        with PILImage.open(p) as im:
            _image_cache[k] = im.size[0] / im.size[1]
    return _image_cache[k]

def _img(slide, path, left, top, width, height=None):
    p = Path(path)
    if not p.exists(): return None
    ar = _ar(p)
    if ar is None: return None
    if height is None: h = width / ar
    elif width is None: width = height * ar; h = height
    else:
        if width / ar <= height: h = width / ar
        else: width = height * ar; h = height
    return slide.shapes.add_picture(str(p), left, top, width, h)

def _tb(slide, l, t, w, h, text, fs=Pt(14), bold=False, color=DARK, align=PP_ALIGN.LEFT):
    tx = slide.shapes.add_textbox(l, t, w, h)
    tx.text_frame.word_wrap = True
    p = tx.text_frame.paragraphs[0]
    p.alignment = align; p.text = text
    p.font.size = fs; p.font.bold = bold; p.font.color.rgb = color; p.font.name = "Arial"
    return tx

def _ml(slide, l, t, w, h, lines, fs=Pt(12), color=DARK):
    tx = slide.shapes.add_textbox(l, t, w, h)
    tf = tx.text_frame; tf.word_wrap = True
    for i, (text, bold, fsize, c) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text; p.font.size = fsize if fsize else fs
        p.font.bold = bold; p.font.color.rgb = c if c else color
        p.font.name = "Arial"; p.space_after = Pt(6)
    return tx

def _pn(s, n):
    _tb(s, Inches(12.3), Inches(7.05), Inches(0.8), Inches(0.35),
        str(n), fs=Pt(10), color=GRAY, align=PP_ALIGN.RIGHT)

def _header(s, n, title, subtitle):
    bar = s.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, Inches(0.07))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
    _tb(s, Inches(0.8), Inches(0.25), Inches(11.5), Inches(0.55),
        title, fs=Pt(28), bold=True, color=DARK)
    _tb(s, Inches(0.8), Inches(0.85), Inches(11.5), Inches(0.45),
        subtitle, fs=Pt(12), color=GRAY)
    _pn(s, n)

# ============================================================
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]

# ---- Slide 1: Cover ----
s = prs.slides.add_slide(blank)
bg = s.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x1A); bg.line.fill.background()
bar = s.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, Inches(0.06))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()

_tb(s, Inches(1.2), Inches(1.5), Inches(10.9), Inches(1.0),
   "YBCO KID 五谐振器微波-光学联合表征", fs=Pt(38), bold=True, color=WHITE)
_ml(s, Inches(1.2), Inches(3.0), Inches(10.9), Inches(2.8), [
    ("样品：YBCO KID 谐振器阵列  |  温度范围：6 K → 76 K", False, Pt(16), LGRAY),
    ("", False, Pt(8), LGRAY),
    ("R1–R5：3.846 / 4.010 / 4.500 / 4.997 / 5.252 GHz", False, Pt(15), LGRAY),
    ("VNA 功率：-25 / -30 / -45 dBm  |  激光功率：0, 1, 3, 5, 7, 9 mW", False, Pt(14), LGRAY),
    ("", False, Pt(8), LGRAY),
    ("组会内部简报  ·  2026-06-18  ·  cmplxIQ 拟合分析", False, Pt(14), GRAY),
], fs=Pt(14), color=LGRAY)
bot = s.shapes.add_shape(1, Inches(1.2), Inches(6.2), Inches(3.0), Inches(0.03))
bot.fill.solid(); bot.fill.fore_color.rgb = ACCENT; bot.line.fill.background()
for i in range(5):
    sw = s.shapes.add_shape(1, Inches(5.5 + i * 0.55), Inches(4.6), Inches(0.45), Inches(0.06))
    sw.fill.solid(); sw.fill.fore_color.rgb = RCOLORS[i]; sw.line.fill.background()
    _tb(s, Inches(5.5 + i * 0.55), Inches(4.7), Inches(0.45), Inches(0.2),
        RESONATORS[i]["name"], fs=Pt(8), color=LGRAY, align=PP_ALIGN.CENTER)
print("[OK] 1 封面")

# ---- Slide 2: Spectrum Overview ----
s = prs.slides.add_slide(blank)
_header(s, 2, "谐振频谱总览",
        "6 K, -25 dBm, 暗态。幅度谷 + 相位差分峰联合判据检出 5 个谐振峰")
_img(s, DETECT / "resonance_detection.jpg", Inches(0.8), Inches(1.5), Inches(11.7))
info = "  |  ".join([f"{r['name']}: {r['freq']}  dip={r['dip']} dB" for r in RESONATORS])
_tb(s, Inches(0.8), Inches(6.7), Inches(11.7), Inches(0.3), info, fs=Pt(10), color=GRAY, align=PP_ALIGN.CENTER)
print("[OK] 2 频谱总览")

# ---- Slide 3: f0(T) + Normalized ----
s = prs.slides.add_slide(blank)
_header(s, 3, "f₀(T) — 谐振频率温度响应",
        "五谐振器 f₀ 均随温度升高单调蓝移，符合超导动能电感 Lₖ ∝ λ²(T)")
_img(s, CMP / "f0_vs_temp_all.jpg", Inches(0.5), Inches(1.5), Inches(6.2))
_img(s, CMP / "normalized_f0_vs_temp.jpg", Inches(6.9), Inches(1.5), Inches(6.2))
print("[OK] 3 f0(T)")

# ---- Slide 4: Qi(T) + Dip vs T ----
s = prs.slides.add_slide(blank)
_header(s, 4, "Qi(T)  &  谐振深度 — 温度响应",
        "左：Qi 随温度升高下降（准粒子热激发损耗增大） |  右：R4 全温最深 dip")
_img(s, CMP / "qi_vs_temp_all.jpg", Inches(0.5), Inches(1.5), Inches(6.2))
_img(s, CMP / "dip_vs_temp.jpg", Inches(6.9), Inches(1.5), Inches(6.2))
print("[OK] 4 Qi+Dip")

# ---- Slides 5-8: Optical Response Comparison @ 4 temps ----
OPT_INFO = [
    ("6K", "深度超导态，准粒子密度极低，光响应呈良好线性"),
    ("20K", "中低温，热准粒子开始贡献，响应率相比 6K 略有下降"),
    ("40K", "中高温，各谐振器响应率分化明显"),
    ("77K", "接近 Tc，超导序参量显著减弱，响应率大幅下降"),
]
for idx, (tlabel, tdesc) in enumerate(OPT_INFO):
    s = prs.slides.add_slide(blank)
    img_path = CMP / f"optical_response_{tlabel}.jpg"
    caption_line = f"五谐振器光学响应对比 — {tlabel}  |  {tdesc}"
    if img_path.exists():
        _img(s, img_path, Inches(0.5), Inches(0.5), Inches(12.3))
        _tb(s, Inches(0.8), Inches(7.0), Inches(11.5), Inches(0.35),
            caption_line, fs=Pt(12), color=DARK)
        _pn(s, 5 + idx)
        print(f"[OK] {5+idx} 光响应 {tlabel}")
    else:
        print(f"[SKIP] {5+idx} 光响应 {tlabel} — 无图")
        _tb(s, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.6),
            f"光学响应 — {tlabel}", fs=Pt(28), bold=True, color=DARK)
        _pn(s, 5 + idx)

# ---- Slide 9: Responsivity vs T ----
s = prs.slides.add_slide(blank)
_header(s, 9, "光学响应率 — 温度依赖性",
        "各谐振器 df₀/dP (kHz/mW) 随温度的变化趋势，反映 Δ(T) 对光生准粒子的调控")
_img(s, CMP / "responsivity_vs_temp.jpg", Inches(1.5), Inches(1.6), Inches(10.3))
print("[OK] 9 响应率")

# ---- Slides 10-14: Per-Resonator Overview ----
for i, (r, rdir) in enumerate(zip(RESONATORS, R_DIRS)):
    s = prs.slides.add_slide(blank)
    sn = 10 + i
    c = RCOLORS[i]
    features = []
    if i == 3: features.append("★ 最深 dip (-17.74 dB)")
    elif i == 1: features.append("最浅 dip (-4.33 dB)")
    _header(s, sn, f"{r['name']} — {r['freq']}  {(' | ' + ' '.join(features)) if features else ''}",
            "左列：f₀(T) + Qi(T)  |  右列：S21 温度叠加 + 响应率")

    # 色标
    sw = s.shapes.add_shape(1, Inches(0.8), Inches(1.55), Inches(0.7), Inches(0.07))
    sw.fill.solid(); sw.fill.fore_color.rgb = c; sw.line.fill.background()

    # 左列
    _img(s, rdir / "f0_versus_temp.jpg", Inches(0.5), Inches(1.8), Inches(5.9))
    _img(s, rdir / "qis_versus_temp.jpg", Inches(0.5), Inches(4.1), Inches(5.9))
    # 右列
    _img(s, rdir / "s21 vs - temp.jpg", Inches(6.8), Inches(1.8), Inches(6.0))
    _img(s, rdir / "responsivity_vs_temp.jpg", Inches(6.8), Inches(4.1), Inches(6.0))

    print(f"[OK] {sn} {r['name']} 概览")

# ---- Slides 15-16: R4 Detail Pages ----
# R4 @ 6K and 40K — use scrapps per-temp images
rdir4 = R_DIRS[3]
for j, (ts, tlbl) in enumerate(zip(TEMP_POINTS[:3], TEMP_LABELS[:3])):  # 6K, 20K, 40K
    # Check if images exist for this resonator at this temperature
    s21_25 = rdir4 / f"s21 - {ts}-25dBm.jpg"
    s21_30 = rdir4 / f"s21 - {ts}-30dBm.jpg"
    s21_45 = rdir4 / f"s21 - {ts}-45dBm.jpg"
    res_shift = rdir4 / f"res shift - {ts}.jpg"

    if s21_25.exists() and res_shift.exists():
        s = prs.slides.add_slide(blank)
        sn = 15 + j
        _header(s, sn, f"R4 (4.997 GHz) 光学响应详情 — {tlbl}",
                f"最深 dip 谐振器 @ {tlbl}  |  S21 × 3 VNA 功率 + 谐振频移 vs 激光功率")

        # 2×2 grid
        imgs = [(s21_25, "S21 @ -25 dBm"), (s21_30, "S21 @ -30 dBm"),
                (s21_45, "S21 @ -45 dBm"), (res_shift, "Res shift vs Laser")]
        pos = [(Inches(0.5), Inches(1.5)), (Inches(6.8), Inches(1.5)),
               (Inches(0.5), Inches(4.3)), (Inches(6.8), Inches(4.3))]
        for (ipath, lbl), (x, y) in zip(imgs, pos):
            _img(s, ipath, x, y, Inches(5.8))
            _tb(s, x, y + Inches(2.55), Inches(5.8), Inches(0.25),
                lbl, fs=Pt(10), color=GRAY, align=PP_ALIGN.CENTER)
        print(f"[OK] {sn} R4 详情 {tlbl}")
    else:
        print(f"[SKIP] R4 详情 {tlbl} — images missing at {ts}")

# ---- Slide 18: R2 Comparison Detail ----
rdir2 = R_DIRS[1]
ts_6k = "5.991K"
ts_77k = "76.204K"
if (rdir2 / f"s21 - {ts_77k}-25dBm.jpg").exists():
    s = prs.slides.add_slide(blank)
    sn = 18
    _header(s, sn, "R2 (4.010 GHz) — 6K vs 77K 光学响应对比",
            "最浅 dip 谐振器在低温与高温下的光响应差异")

    # Top: 6K res shift + 77K res shift side by side
    _img(s, rdir2 / f"res shift - {ts_6k}.jpg", Inches(0.5), Inches(1.5), Inches(5.8))
    _tb(s, Inches(0.5), Inches(4.1), Inches(5.8), Inches(0.25),
        "Res shift @ 6K", fs=Pt(10), color=GRAY, align=PP_ALIGN.CENTER)
    _img(s, rdir2 / f"res shift - {ts_77k}.jpg", Inches(6.8), Inches(1.5), Inches(5.8))
    _tb(s, Inches(6.8), Inches(4.1), Inches(5.8), Inches(0.25),
        "Res shift @ 77K", fs=Pt(10), color=GRAY, align=PP_ALIGN.CENTER)

    # Bottom: S21 snapshots
    _img(s, rdir2 / f"s21 - {ts_6k}-25dBm.jpg", Inches(0.5), Inches(4.55), Inches(3.8))
    _img(s, rdir2 / f"s21 - {ts_77k}-25dBm.jpg", Inches(4.5), Inches(4.55), Inches(3.8))
    _tb(s, Inches(0.5), Inches(6.95), Inches(3.8), Inches(0.25),
        "S21 -25dBm @ 6K", fs=Pt(9), color=GRAY, align=PP_ALIGN.CENTER)
    _tb(s, Inches(4.5), Inches(6.95), Inches(3.8), Inches(0.25),
        "S21 -25dBm @ 77K", fs=Pt(9), color=GRAY, align=PP_ALIGN.CENTER)
    print(f"[OK] {sn} R2 对比")
else:
    print(f"[SKIP] R2 对比 — 77K images missing")

# ---- Slide 19: R1/R3/R5 Mini Comparison ----
s = prs.slides.add_slide(blank)
_header(s, 19, "R1 / R3 / R5 — 谐振频率温度响应对比",
        "f₀(T) 对比：三个中浅谐振器的蓝移行为一致")
for i in [0, 2, 4]:  # R1, R3, R5
    x = Inches(0.4 + i/2 * 4.2)
    _img(s, R_DIRS[i] / "f0_versus_temp.jpg", x, Inches(1.5), Inches(4.0))
    _tb(s, x, Inches(4.55), Inches(4.0), Inches(0.25),
        f"{RESONATORS[i]['name']} ({RESONATORS[i]['freq']})", fs=Pt(10), color=GRAY, align=PP_ALIGN.CENTER)
print("[OK] 19 R1/R3/R5")

# ---- Slide 20: Summary ----
s = prs.slides.add_slide(blank)
bar = s.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, Inches(0.07))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
_tb(s, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.6),
   "小结与下一步", fs=Pt(32), bold=True, color=DARK)
_ml(s, Inches(1.0), Inches(1.5), Inches(11.3), Inches(5.2), [
    ("五谐振器 cmplxIQ 拟合分析完成", True, Pt(17), DARK),
    ("", False, Pt(7), DARK),
    ("R1–R5：3.846 / 4.010 / 4.500 / 4.997 / 5.252 GHz，全温稳定工作", False, Pt(14), GRAY),
    ("f₀(T) 蓝移一致 — 五谐振器均随温度升高单调蓝移，符合动能电感标度行为", False, Pt(14), GRAY),
    ("Qi(T) 趋势一致 — 低温段高 Qi（~2000–5000），随温度升高逐渐下降", False, Pt(14), GRAY),
    ("★ R4 (4.997 GHz) 性能最优 — dip = -17.74 dB（最深），Qi 最高，推荐作为主读出像素", False, Pt(14), GRAY),
    ("光学响应 — 6K/20K/40K/77K 四个温度点五谐振器响应率系统对比", False, Pt(14), GRAY),
    ("响应率差异 — 存在与频率/耦合相关的系统性差异，R4 响应率值得特别关注", False, Pt(14), GRAY),
    ("", False, Pt(10), DARK),
    ("下一步：R4 优先做 NEP 估算、多像素统计相关性分析、低温 LNA 集成测试", True, Pt(14), ACCENT2),
], fs=Pt(14), color=GRAY)
_pn(s, 20)
print("[OK] 20 小结")

# ============================================================
prs.save(str(PPTX_OUT))
print(f"\n[OK] PPT saved: {PPTX_OUT}")
print(f"     Total {len(prs.slides)} slides")
