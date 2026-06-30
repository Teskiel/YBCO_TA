# -*- coding: utf-8 -*-
"""
生成 YBCO KID 微波-光学联合表征简报 v7。

基于 merged/ 下 8 个分析文件夹（不含 compare），
选取 6K / 20K / 40K / 77K 四个代表温度点各做一页详情。

v7 修订：
- 所有图片等比例缩放（仅指定宽，高度按原图 AR 计算）
- 三张概览图拆为两页（2+1）
- 温度详情页采用 1 大图 + 3 小图 布局
"""

from pathlib import Path
from PIL import Image as PILImage
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# ============================================================
# 路径配置
# ============================================================
SCRIPT_DIR = Path(__file__).resolve().parent
MERGED = SCRIPT_DIR / "output" / "merged"
PPTX_OUT = SCRIPT_DIR / "output" / "YBCO_KID_merged_表征简报_v7.pptx"

F01 = MERGED / "01_resonance_detection"
F02 = MERGED / "02_f0_temperature"
F03 = MERGED / "03_Qi_temperature"
F04 = MERGED / "04_S21_temperature_overlay"
F05 = MERGED / "05_optical_response_6K"
F06 = MERGED / "06_optical_response_highT"
F07 = MERGED / "07_responsivity_temperature"
F08 = MERGED / "08_per_temp_raw"

# ============================================================
# 颜色 / 字体
# ============================================================
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
DARK    = RGBColor(0x1A, 0x1A, 0x1A)
GRAY    = RGBColor(0x66, 0x66, 0x66)
LGRAY   = RGBColor(0xAA, 0xAA, 0xAA)
ACCENT  = RGBColor(0x1F, 0x77, 0xB4)
ACCENT2 = RGBColor(0xD6, 0x27, 0x28)
LIGHT_BG = RGBColor(0xF5, 0xF5, 0xF5)

FONT_TITLE = "Arial"
FONT_BODY  = "Arial"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ============================================================
# 图片缓存（避免反复打开）
# ============================================================
_image_cache = {}

def _get_image_ar(img_path):
    """获取图片宽高比 (w/h)，带缓存。"""
    p = Path(img_path)
    if not p.exists():
        return None
    key = str(p)
    if key not in _image_cache:
        with PILImage.open(p) as im:
            _image_cache[key] = im.size[0] / im.size[1]
    return _image_cache[key]


def _add_image_proportional(slide, img_path, left, top, width, height=None):
    """
    等比例添加图片。
    - 如果只给 width：高度 = width / AR
    - 如果只给 height：宽度 = height * AR
    - 如果都给：取两者中更小的那个方向，另一方向按 AR 计算
    """
    p = Path(img_path)
    if not p.exists():
        print(f"  [WARN] 图片不存在: {p}")
        return None

    ar = _get_image_ar(p)
    if ar is None:
        return None

    if width is not None and height is None:
        # 定宽，算高
        h = width / ar
    elif height is not None and width is None:
        # 定高，算宽
        w = height * ar
    else:
        # 两个都给了，取 fit-in
        if width / ar <= height:
            h = width / ar
        else:
            width = height * ar
            h = height

    return slide.shapes.add_picture(str(p), left, top, width, h)


def _add_textbox(slide, left, top, width, height, text,
                 font_size=Pt(14), bold=False, color=DARK,
                 alignment=PP_ALIGN.LEFT, font_name=FONT_BODY):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    p = txBox.text_frame.paragraphs[0]
    p.alignment = alignment
    p.text = text
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    return txBox


def _add_multiline(slide, left, top, width, height, lines,
                   font_size=Pt(12), color=DARK, font_name=FONT_BODY):
    """lines: [(text, bold, font_size, color), ...]"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, bold, fs, c) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = fs if fs else font_size
        p.font.bold = bold
        p.font.color.rgb = c if c else color
        p.font.name = font_name
        p.space_after = Pt(6)
    return txBox


def _add_page_number(slide, num):
    _add_textbox(slide, Inches(12.3), Inches(7.05), Inches(0.8), Inches(0.35),
                 str(num), font_size=Pt(10), color=GRAY, alignment=PP_ALIGN.RIGHT)


def _add_section_header(slide, num, title, subtitle_text):
    """统一章节页头：顶部色条 + 标题 + 副标题。"""
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, Inches(0.07))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    _add_textbox(slide, Inches(0.8), Inches(0.25), Inches(11.5), Inches(0.55),
                 title, font_size=Pt(28), bold=True, color=DARK)
    _add_textbox(slide, Inches(0.8), Inches(0.85), Inches(11.5), Inches(0.5),
                 subtitle_text, font_size=Pt(12), color=GRAY)
    _add_page_number(slide, num)


# ============================================================
# 构建 PPT
# ============================================================
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]


# ════════════════════════════════════════════════════════════
# Slide 1 — 封面
# ════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank)
bg = s1.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
bg.fill.solid()
bg.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
bg.line.fill.background()

bar = s1.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, Inches(0.06))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT
bar.line.fill.background()

_add_textbox(s1, Inches(1.2), Inches(1.8), Inches(10.9), Inches(1.0),
             "YBCO KID 微波-光学联合表征",
             font_size=Pt(40), bold=True, color=WHITE)

cover_lines = [
    ("样品：YBCO KID 谐振器  |  温度范围：6 K → 76 K（38 个温度点）", False, Pt(16), LGRAY),
    ("", False, Pt(8), LGRAY),
    ("VNA 读出功率：-25 / -30 / -45 dBm  |  激光功率：0, 1, 3, 5, 7, 9 mW", False, Pt(14), LGRAY),
    ("", False, Pt(8), LGRAY),
    ("组会内部简报  ·  2026-06-18", False, Pt(14), GRAY),
]
_add_multiline(s1, Inches(1.2), Inches(3.2), Inches(10.9), Inches(2.5),
               cover_lines, font_size=Pt(14), color=LGRAY)

bot_bar = s1.shapes.add_shape(1, Inches(1.2), Inches(6.2), Inches(3.0), Inches(0.03))
bot_bar.fill.solid()
bot_bar.fill.fore_color.rgb = ACCENT
bot_bar.line.fill.background()
print("  [OK] Slide 1 — 封面")


# ════════════════════════════════════════════════════════════
# Slide 2 — 谐振峰识别
# ════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank)
_add_section_header(s2, 2,
    "谐振峰自动识别",
    "采用幅度谷 + 相位差分峰联合判据自动寻峰（SNR ≥ 0.5），选定谐振峰位于 ~3.X GHz")

img = F01 / "resonance_detection.jpg"
img_w = Inches(10.5)
_add_image_proportional(s2, img, Inches(1.4), Inches(1.5), img_w)
print("  [OK] Slide 2 — 谐振峰识别")


# ════════════════════════════════════════════════════════════
# Slide 3 — 全温趋势 A：S21 叠加 + f₀(T)
# ════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank)
_add_section_header(s3, 3,
    "全温 S21 叠加  &  f₀(T) 温度响应",
    "S21 叠加显示谐振峰随温度单调蓝移 → 超导动能电感效应  |  f₀ 随温度升高单调下降，符合 Lₖ ∝ λ²(T)")

# 两张图并列
img_w = Inches(5.6)
gap = Inches(0.35)
start_x = Inches(0.9)
top_y = Inches(1.6)
labels = ["全温 S21 叠加", "f₀(T) 谐振频率 vs 温度"]

for i, (f, lbl) in enumerate([(F04 / "s21 vs - temp.jpg", labels[0]),
                               (F02 / "f0_versus_temp.jpg", labels[1])]):
    x = start_x + i * (img_w + gap)
    _add_image_proportional(s3, f, x, top_y, img_w)
    _add_textbox(s3, x, top_y + Inches(4.45), img_w, Inches(0.3),
                 lbl, font_size=Pt(11), color=GRAY, alignment=PP_ALIGN.CENTER)

print("  [OK] Slide 3 — 全温趋势 A")


# ════════════════════════════════════════════════════════════
# Slide 4 — 全温趋势 B：Qi(T)
# ════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank)
_add_section_header(s4, 4,
    "Qi(T) 内禀品质因数温度响应",
    "Qi 低温段较高，随温度上升逐渐降低 → 准粒子热激发损耗增大  |  三个 VNA 功率偏差小，读出功率未引入显著非线性")

img_w = Inches(8.5)
img = F03 / "qis_versus_temp.jpg"
ar = _get_image_ar(img)
if ar:
    _add_image_proportional(s4, img,
                            (SLIDE_W - img_w) / 2, Inches(1.6), img_w)
print("  [OK] Slide 4 — 全温趋势 B")


# ════════════════════════════════════════════════════════════
# Slide 5–8 — 四个温度点详情页
# 布局：左侧大图 (res shift) + 右侧三张小图 (S21 × 3 功率) 纵向堆叠
# ════════════════════════════════════════════════════════════

TEMP_POINTS = [
    {
        "label": "6 K",
        "temp_str": "5.991K",
        "title": "低温光致响应 — 6 K",
        "subtitle": ("深度超导态，准粒子密度极低。"
                     "光注入非平衡准粒子 → 破坏库珀对 → 动能电感增大 → 谐振频率红移。"
                     "三个 VNA 功率下响应率一致，表明读出功率未加热器件。"),
        "s21_dir": F05,
        "res_shift_dir": F05,
    },
    {
        "label": "20 K",
        "temp_str": "19.977K",
        "title": "中低温光致响应 — 20 K",
        "subtitle": ("超导能隙仍较大，热准粒子开始贡献。"
                     "相比 6 K，本底准粒子密度升高 → 光注入的相对增量减小 → 响应率略有下降。"
                     "频移-激光功率仍保持良好线性。"),
        "s21_dir": F08,
        "res_shift_dir": F07,
    },
    {
        "label": "40 K",
        "temp_str": "39.825K",
        "title": "中高温光致响应 — 40 K",
        "subtitle": ("接近超导转变中段 (T/Tc ~ 0.5)，动能电感温度敏感性增强。"
                     "热准粒子密度显著升高，光致频移量级明显减小。"
                     "S21 峰形仍保持良好，器件在此温区仍稳定工作。"),
        "s21_dir": F08,
        "res_shift_dir": F07,
    },
    {
        "label": "77 K",
        "temp_str": "76.204K",
        "title": "高温光致响应 — 77 K",
        "subtitle": ("接近 Tc (~85-90 K)，超导序参量 Δ(T) 显著减弱。"
                     "热准粒子主导，光学响应率相比 6 K 大幅下降。"
                     "谐振峰仍可分辨，但 Qi 明显降低。"),
        "s21_dir": F06,
        "res_shift_dir": F07,
    },
]

# 大图宽度
LARGE_W = Inches(5.5)
# 小图宽度（三张纵向排列在右侧）
SMALL_W = Inches(3.6)
# 小图之间的纵向间距
SMALL_GAP_Y = Inches(0.15)

for idx, tp in enumerate(TEMP_POINTS):
    s = prs.slides.add_slide(blank)
    slide_num = 5 + idx
    _add_section_header(s, slide_num, tp["title"], tp["subtitle"])

    ts = tp["temp_str"]
    s21_dir = tp["s21_dir"]
    res_dir = tp["res_shift_dir"]

    # --- 大图：谐振频移 vs 激光功率 ---
    large_img = res_dir / f"res shift - {ts}.jpg"
    large_top = Inches(1.6)
    large_left = Inches(0.6)
    _add_image_proportional(s, large_img, large_left, large_top, LARGE_W)
    # 大图标签
    large_ar = _get_image_ar(large_img)
    large_h = LARGE_W / large_ar if large_ar else Inches(4.3)
    _add_textbox(s, large_left, large_top + large_h + Inches(0.05), LARGE_W, Inches(0.28),
                 f"谐振频移 vs 激光功率  @ {tp['label']}",
                 font_size=Pt(10), color=GRAY, alignment=PP_ALIGN.CENTER)

    # --- 三张小图：S21 × 三个 VNA 功率 ---
    small_imgs = [
        (s21_dir / f"s21 - {ts}-25dBm.jpg", "S21  @ -25 dBm"),
        (s21_dir / f"s21 - {ts}-30dBm.jpg", "S21  @ -30 dBm"),
        (s21_dir / f"s21 - {ts}-45dBm.jpg", "S21  @ -45 dBm"),
    ]

    # 右侧小图起始位置
    small_left = Inches(7.2)
    small_top_start = Inches(1.6)

    for j, (simg, slbl) in enumerate(small_imgs):
        s_top = small_top_start + j * (SMALL_GAP_Y + Inches(1.8))
        # 检查图片是否存在
        sar = _get_image_ar(simg)
        if sar is None:
            print(f"  [WARN] missing: {simg}")
            continue
        s_h = SMALL_W / sar
        _add_image_proportional(s, simg, small_left, s_top, SMALL_W)
        _add_textbox(s, small_left, s_top + s_h + Inches(0.03), SMALL_W, Inches(0.25),
                     slbl, font_size=Pt(9), color=GRAY, alignment=PP_ALIGN.CENTER)

    print(f"  [OK] Slide {slide_num} — {tp['label']} 详情")


# ════════════════════════════════════════════════════════════
# Slide 9 — 响应率 vs 温度汇总
# ════════════════════════════════════════════════════════════
s9 = prs.slides.add_slide(blank)
_add_section_header(s9, 9,
    "光学响应率 — 温度依赖性",
    "从 6 K 到 76 K 各温度点的谐振频移 vs 激光功率拟合斜率汇总  |  响应率 (Hz/W) 趋势反映 Δ(T) 对光生准粒子的调控")

img = F07 / "responsivity_vs_temp.jpg"
img_w = Inches(10.0)
ar = _get_image_ar(img)
if ar:
    _add_image_proportional(s9, img,
                            (SLIDE_W - img_w) / 2, Inches(1.5), img_w)
print("  [OK] Slide 9 — 响应率 vs 温度")


# ════════════════════════════════════════════════════════════
# Slide 10 — 小结
# ════════════════════════════════════════════════════════════
s10 = prs.slides.add_slide(blank)
bar = s10.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, Inches(0.07))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT
bar.line.fill.background()

_add_textbox(s10, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.6),
             "小结与下一步", font_size=Pt(32), bold=True, color=DARK)

summary = [
    ("成功表征 YBCO KID 在 6–76 K 的微波谐振特性与光学响应", True, Pt(16), DARK),
    ("", False, Pt(6), DARK),
    ("谐振峰识别 — 幅度谷 + 相位差分峰联合判据自动寻峰，器件全温稳定工作", False, Pt(14), GRAY),
    ("f₀(T) 蓝移 — 随温度升高单调蓝移，符合超导动能电感 Lₖ ∝ λ²(T) 理论预期", False, Pt(14), GRAY),
    ("Qi(T) 下降 — 随温度上升逐渐降低，归因于准粒子热激发损耗增大", False, Pt(14), GRAY),
    ("光学响应 — 6K / 20K / 40K / 77K 四个代表温度点频移 vs 激光功率均呈良好线性", False, Pt(14), GRAY),
    ("响应率温度依赖 — 随温度单调变化，与超导能隙 Δ(T) 定性一致", False, Pt(14), GRAY),
    ("", False, Pt(10), DARK),
    ("下一步：NEP 噪声等效功率估算、多像素统计比较、低温放大器集成测试", True, Pt(14), ACCENT2),
]
_add_multiline(s10, Inches(1.0), Inches(1.5), Inches(11.3), Inches(5.0),
               summary, font_size=Pt(14), color=GRAY)
_add_page_number(s10, 10)
print("  [OK] Slide 10 — 小结")


# ============================================================
# 保存
# ============================================================
prs.save(str(PPTX_OUT))
print(f"\n[OK] PPT saved: {PPTX_OUT}")
print(f"     Total {len(prs.slides)} slides")
