# -*- coding: utf-8 -*-
"""
YBCO KID 新旧测量对比简报
对比 accomplish_merged (OLD, 20260612-0614) vs 20260609-0624 (NEW)
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image as PILImage

SCRIPT_DIR = Path(__file__).resolve().parent
OUT_DIR = SCRIPT_DIR / "output" / "comparison"
IMG_DIR = OUT_DIR  # comparison plots are here
PLOT_DIR = Path("../Auto_Sweep/experiment_data/~merged/output/_cache/plot_output")
PPTX_OUT = OUT_DIR / "YBCO_KID_新旧测量对比简报_v2.pptx"

# Colors
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x1A, 0x1A, 0x1A)
GRAY   = RGBColor(0x66, 0x66, 0x66)
LGRAY  = RGBColor(0xAA, 0xAA, 0xAA)
ACCENT = RGBColor(0x1F, 0x77, 0xB4)
ACCENT2 = RGBColor(0xD6, 0x27, 0x28)
GREEN  = RGBColor(0x2C, 0xA0, 0x2C)
ORANGE = RGBColor(0xFF, 0x7F, 0x0E)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

_image_cache = {}
def _ar(p):
    p = Path(p)
    if not p.exists(): return None
    k = str(p)
    if k not in _image_cache:
        with PILImage.open(p) as im:
            _image_cache[k] = im.size[0] / im.size[1]
    return _image_cache[k]

def _img(slide, path, left, top, width, height=None):
    p = Path(path)
    if not p.exists(): return None
    ar = _ar(p)
    if ar is None: return None
    if height is None: h = width / ar
    elif width is None: width = height * ar; h = height
    else:
        if width / ar <= height: h = width / ar
        else: width = height * ar; h = height
    return slide.shapes.add_picture(str(p), left, top, width, h)

def _tb(slide, l, t, w, h, text, fs=Pt(14), bold=False, color=DARK, align=PP_ALIGN.LEFT):
    tx = slide.shapes.add_textbox(l, t, w, h)
    tx.text_frame.word_wrap = True
    p = tx.text_frame.paragraphs[0]
    p.alignment = align; p.text = text
    p.font.size = fs; p.font.bold = bold; p.font.color.rgb = color; p.font.name = "Arial"
    return tx

def _ml(slide, l, t, w, h, lines, fs=Pt(12), color=DARK):
    tx = slide.shapes.add_textbox(l, t, w, h)
    tf = tx.text_frame; tf.word_wrap = True
    for i, (text, bold, fsize, c) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text; p.font.size = fsize if fsize else fs
        p.font.bold = bold; p.font.color.rgb = c if c else color
        p.font.name = "Arial"; p.space_after = Pt(6)
    return tx

def _pn(s, n):
    _tb(s, Inches(12.3), Inches(7.05), Inches(0.8), Inches(0.35),
        str(n), fs=Pt(10), color=GRAY, align=PP_ALIGN.RIGHT)

def _add_blank_slide():
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)

def _add_bg(slide, color=DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def _add_dark_header(slide, title, subtitle=""):
    _tb(slide, Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.7),
        title, fs=Pt(32), bold=True, color=WHITE)
    if subtitle:
        _tb(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.5),
            subtitle, fs=Pt(14), color=LGRAY)
    # accent line
    line = slide.shapes.add_shape(1, Inches(0.8), Inches(1.85), Inches(2.5), Pt(3))
    line.fill.solid(); line.fill.fore_color.rgb = ACCENT; line.line.fill.background()

# ============================================================
# SLIDE 1: Title
# ============================================================
s = _add_blank_slide()
_add_bg(s, DARK)
_tb(s, Inches(1.0), Inches(1.5), Inches(11.0), Inches(1.2),
    "YBCO KID 五谐振器表征", fs=Pt(40), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
_tb(s, Inches(1.0), Inches(2.8), Inches(11.0), Inches(0.8),
    "新旧两次测量数据对比分析", fs=Pt(24), color=ACCENT, align=PP_ALIGN.CENTER)
line = s.shapes.add_shape(1, Inches(4.5), Inches(3.8), Inches(4.0), Pt(2))
line.fill.solid(); line.fill.fore_color.rgb = ACCENT; line.line.fill.background()
_tb(s, Inches(1.0), Inches(4.2), Inches(11.0), Inches(0.5),
    "OLD: accomplish_merged (2026-06-12 ~ 06-14)  |  NEW: 20260609-0624 (06-09 ~ 06-24)",
    fs=Pt(14), color=LGRAY, align=PP_ALIGN.CENTER)
_tb(s, Inches(1.0), Inches(5.0), Inches(11.0), Inches(0.5),
    "2026-06-25", fs=Pt(12), color=GRAY, align=PP_ALIGN.CENTER)
_pn(s, 1)

# ============================================================
# SLIDE 2: Experiment Parameters Comparison
# ============================================================
s = _add_blank_slide()
_add_bg(s, DARK)
_add_dark_header(s, "实验参数对比", "OLD (accomplish_merged) vs NEW (20260609-0624)")

param_lines = [
    ("参数", True, Pt(14), WHITE),
    ("", False, Pt(6), WHITE),
    ("温度范围", True, Pt(14), ACCENT),
    ("  OLD: 6 ~ 80 K, 步长 2 K, 共 38 个温度点", False, Pt(13), WHITE),
    ("  NEW: 6, 10, 20, 40, 50, 60, 70, 77 K, 共 8 个温度点", False, Pt(13), WHITE),
    ("", False, Pt(6), WHITE),
    ("VNA 功率", True, Pt(14), ACCENT),
    ("  OLD: -25, -30, -45 dBm (3 级)", False, Pt(13), WHITE),
    ("  NEW: -55 ~ -25 dBm, 步长 2 dB, 共 16 级", False, Pt(13), WHITE),
    ("", False, Pt(6), WHITE),
    ("激光功率", True, Pt(14), ACCENT),
    ("  OLD: 0, 1, 3, 5, 7, 9 mW (6 级)", False, Pt(13), WHITE),
    ("  NEW: 0, 1, 3, 5, 7, 9 mW (6 级, T<=40K); 扩展至 17 mW (T>=50K)", False, Pt(13), WHITE),
    ("", False, Pt(6), WHITE),
    ("数据量", True, Pt(14), ACCENT),
    ("  OLD: 668 S2P 文件  |  NEW: ~960 S2P 文件 (不含 80K)", False, Pt(13), WHITE),
    ("", False, Pt(6), WHITE),
    ("共同覆盖温度", True, Pt(14), ACCENT),
    ("  6, 10, 20, 40, 50, 60, 70 K — 7 个共同温度点可用于直接对比", False, Pt(13), WHITE),
]
_ml(s, Inches(0.8), Inches(2.2), Inches(11.5), Inches(5.0), param_lines)
_pn(s, 2)

# ============================================================
# SLIDE 3: Temperature Stability Summary
# ============================================================
s = _add_blank_slide()
_add_bg(s, DARK)
_add_dark_header(s, "各温度数据稳定性评级 (NEW 数据集)")

# Table: temperature | resonators | avg dip depth | VNA levels | Laser levels | Grade
stab_data = [
    ("6 K",  "5/5", "14.5 dB", "16", "6", "Excellent", "A"),
    ("10 K", "5/5", "14.0 dB", "16", "6", "Excellent", "A"),
    ("20 K", "5/5", "12.9 dB", "16", "6", "Excellent", "A"),
    ("40 K", "5/5", "10.6 dB", "16", "6", "Excellent", "A"),
    ("50 K", "5/5", "11.9 dB", "16", "10", "Good", "A"),
    ("60 K", "5/5", "11.7 dB", "16", "10", "Good", "A"),
    ("70 K", "5/5", "3.6 dB", "16", "10", "R1/R2 corrected", "B"),
    ("77 K", "5/5", "3.7 dB", "16", "6", "Dips shallow", "B"),
]

y0 = 2.3
row_h = 0.55
cols = [(0.8, 1.2), (2.2, 0.8), (3.2, 1.2), (4.6, 1.0), (5.8, 1.2), (7.2, 2.0), (9.5, 0.8)]
headers = ["Temperature", "Resonators", "Avg Dip Depth", "VNA Lvls", "Laser Lvls", "Notes", "Grade"]
for j, (hdr, (cx, cw)) in enumerate(zip(headers, cols)):
    _tb(s, Inches(cx), Inches(y0 - 0.4), Inches(cw), Inches(0.35),
        hdr, fs=Pt(11), bold=True, color=LGRAY)

for i, (temp, res, dip, vna, laser, note, grade_color) in enumerate(stab_data):
    y = y0 + (i + 1) * row_h
    _tb(s, Inches(cols[0][0]), Inches(y), Inches(cols[0][1]), Inches(0.35), temp, fs=Pt(12), bold=True, color=WHITE)
    _tb(s, Inches(cols[1][0]), Inches(y), Inches(cols[1][1]), Inches(0.35), res, fs=Pt(12), color=WHITE)
    _tb(s, Inches(cols[2][0]), Inches(y), Inches(cols[2][1]), Inches(0.35), dip, fs=Pt(12), color=WHITE)
    _tb(s, Inches(cols[3][0]), Inches(y), Inches(cols[3][1]), Inches(0.35), vna, fs=Pt(12), color=WHITE)
    _tb(s, Inches(cols[4][0]), Inches(y), Inches(cols[4][1]), Inches(0.35), laser, fs=Pt(12), color=WHITE)
    _tb(s, Inches(cols[5][0]), Inches(y), Inches(cols[5][1]), Inches(0.35), note, fs=Pt(10), color=LGRAY)
    _tb(s, Inches(cols[6][0]), Inches(y), Inches(cols[6][1]), Inches(0.35), grade_color, fs=Pt(14), bold=True, color=ACCENT)

# Notes
_ml(s, Inches(0.8), Inches(y0 + 9 * row_h + 0.2), Inches(11.5), Inches(1.2), [
    ("Grade A: 5/5 resonators, avg dip > 5 dB, excellent for analysis", False, Pt(11), LGRAY),
    ("Grade B: 5/5 resonators detected but dips shallow (< 5 dB), usable with caution", False, Pt(11), LGRAY),
    ("Note: 70K R2 identification corrected from algorithm default (verified at 3.719 GHz)", False, Pt(11), ORANGE),
])
_pn(s, 3)

# ============================================================
# SLIDE 4: f0(T) Comparison Plot
# ============================================================
s = _add_blank_slide()
_add_bg(s, DARK)
_add_dark_header(s, "f0(T) 对比: OLD vs NEW", "所有 5 个谐振子，重叠温度点 6-70K")

_img(s, IMG_DIR / "f0_comparison_all.png",
     Inches(0.5), Inches(2.1), Inches(12.3))
_pn(s, 4)

# ============================================================
# SLIDE 5: Delta f0 Analysis
# ============================================================
s = _add_blank_slide()
_add_bg(s, DARK)
_add_dark_header(s, "频率偏移分析: Delta f0 = OLD - NEW", "系统性偏差随温度升高而增大")

_img(s, IMG_DIR / "delta_f0_vs_T.png",
     Inches(0.5), Inches(2.1), Inches(7.5))

# Summary box
delta_lines = [
    ("偏移量级", True, Pt(14), ACCENT),
    ("", False, Pt(6), WHITE),
    ("6-20 K: < 0.3 MHz", False, Pt(13), WHITE),
    ("  — 两次测量在低温段完全一致", False, Pt(11), GREEN),
    ("40-60 K: 3 ~ 11 MHz 系统性偏移 (OLD > NEW)", False, Pt(13), WHITE),
    ("  — 对应约 2-5 K 的等效温度差异", False, Pt(11), LGRAY),
    ("70 K: 20 ~ 25 MHz (排除 R2 异常值)", False, Pt(13), WHITE),
    ("  — 对应约 8-10 K 等效温差", False, Pt(11), LGRAY),
    ("", False, Pt(8), WHITE),
    ("OLD 测量频率系统偏高 = OLD 样品温度偏低", False, Pt(12), ORANGE),
    ("可能原因: LakeShore 温度计校准漂移或热锚差异", False, Pt(12), ORANGE),
]
_ml(s, Inches(8.5), Inches(2.2), Inches(4.5), Inches(5.0), delta_lines)
_pn(s, 5)

# ============================================================
# SLIDE 6: 70K 区域详细分析 — Dip Depth & df/dT
# ============================================================
s = _add_blank_slide()
_add_bg(s, DARK)
_add_dark_header(s, "70K 区域反常趋势分析: 60K → 70K → 77K", "Dip depth 与频率红移率 df/dT 的逐谐振子对比 (-55dBm, 00mW)")

# Table 1: Dip depth
_tb(s, Inches(0.8), Inches(2.1), Inches(5.5), Inches(0.35),
    "Dip Depth (dB) vs Temperature", fs=Pt(13), bold=True, color=ACCENT)

dip_table = [
    ("R1", "8.03", "4.00", "4.38", "-50%", "+10%", "70K dip 降至 4dB 后趋稳"),
    ("R2", "5.67", "5.16", "3.83", "-9%", "-26%", "70K 仍较深, 77K 加速退化"),
    ("R3", "7.41", "3.83", "3.70", "-48%", "-3%", "70->77 几乎持平"),
    ("R4", "12.45", "8.92", "4.50", "-28%", "-50%", "77K 退化加速"),
    ("R5", "26.46", "4.39", "1.97", "-83%", "-55%", "60->70 暴跌 6 倍!"),
]

dip_cols = [(0.8,0.6), (1.6,0.7), (2.5,0.7), (3.4,0.7), (4.3,0.7), (5.2,0.7), (6.2, 3.8)]
dip_hdrs = ["R", "60K", "70K", "77K", "60->70", "70->77", "Notes"]
for j, (hdr, (cx, cw)) in enumerate(zip(dip_hdrs, dip_cols)):
    _tb(s, Inches(cx), Inches(2.5), Inches(cw), Inches(0.3), hdr, fs=Pt(9), bold=True, color=LGRAY)

for i, row in enumerate(dip_table):
    y = 2.9 + (i) * 0.42
    for j, (val, (cx, cw)) in enumerate(zip(row, dip_cols)):
        color = ORANGE if i == 4 else WHITE  # highlight R5
        bold = True if i == 4 else False
        _tb(s, Inches(cx), Inches(y), Inches(cw), Inches(0.3), val, fs=Pt(10), bold=bold, color=color)

# Table 2: df/dT
_tb(s, Inches(0.8), Inches(5.1), Inches(5.5), Inches(0.35),
    "df/dT (MHz/K) vs Temperature", fs=Pt(13), bold=True, color=ACCENT)

dfdt_table = [
    ("R1", "-12.3", "-25.0", "2.03x", "红移加速, 仍在正常范围"),
    ("R2", "-13.1", "-26.1", "1.99x", "与 R1 同步加速, 行为一致"),
    ("R3", "-13.9", "-28.7", "2.07x", "典型近 Tc 行为"),
    ("R4", "-15.5", "-30.4", "1.96x", "典型近 Tc 行为"),
    ("R5", "-16.4", "-33.3", "2.03x", "df/dT 翻倍, 同时 dip 暴跌"),
]

dfdt_cols = [(0.8,0.6), (1.6,0.9), (2.7,0.9), (3.8,0.9), (6.2, 4.0)]
dfdt_hdrs = ["R", "60-70K", "70-77K", "Ratio", "Notes"]
for j, (hdr, (cx, cw)) in enumerate(zip(dfdt_hdrs, dfdt_cols)):
    _tb(s, Inches(cx), Inches(5.5), Inches(cw), Inches(0.3), hdr, fs=Pt(9), bold=True, color=LGRAY)

for i, row in enumerate(dfdt_table):
    y = 5.9 + (i) * 0.42
    for j, (val, (cx, cw)) in enumerate(zip(row, dfdt_cols)):
        color = ORANGE if i == 1 else WHITE
        bold = True if i == 1 else False
        _tb(s, Inches(cx), Inches(y), Inches(cw), Inches(0.3), val, fs=Pt(10), bold=bold, color=color)

# Key takeaway
_ml(s, Inches(7.5), Inches(2.1), Inches(5.5), Inches(5.0), [
    ("关键观察 (修正 70K R1/R2 互换后)", True, Pt(14), ACCENT),
    ("", False, Pt(6), WHITE),
    ("R5 在 60→70K 间 dip 暴跌 83%", False, Pt(12), WHITE),
    ("  — 从最深 (26.5 dB) 跌至平均水平 (4.4 dB)", False, Pt(10), LGRAY),
    ("  — 不是数据错误: 77K 进一步降至 2.0 dB,", False, Pt(10), LGRAY),
    ("    趋势连续, 物理自洽", False, Pt(10), LGRAY),
    ("", False, Pt(6), WHITE),
    ("所有五个谐振子 df/dT 均匀翻倍 (~2.0x)", False, Pt(12), WHITE),
    ("  — 修正后 R1-R5 的 70→77K 红移率一致加速", False, Pt(10), LGRAY),
    ("  — df/dT 从 ~13-16 升至 ~25-33 MHz/K", False, Pt(10), LGRAY),
    ("  — 无个别异常, 行为高度一致", False, Pt(10), LGRAY),
    ("", False, Pt(6), WHITE),
    ("70K 是拐点而非异常", False, Pt(12), WHITE),
    ("  — 60K 时所有谐振子处于 BCS 平台区", False, Pt(10), LGRAY),
    ("  — 70K 时全部开始进入近 Tc 退化区", False, Pt(10), LGRAY),
    ("  — 77K 时退化加深但五个谐振子步调一致", False, Pt(10), LGRAY),
    ("", False, Pt(6), WHITE),
    ("修正内容: 算法在 70K 互换了 R1 和 R2 的 f0", False, Pt(10), ORANGE),
    ("(R1 正确=3.573, R2 正确=3.719; 77K 无需修正)", False, Pt(10), ORANGE),
    ("实测温度: 60K→60.29K, 70K→70.41K, 77K→76.88K", False, Pt(10), GREEN),
])
_pn(s, 6)

# ============================================================
# SLIDE 7: Physical Explanation — Two-Fluid Model
# ============================================================
s = _add_blank_slide()
_add_bg(s, DARK)
_add_dark_header(s, "物理解释: 两流体模型与近 Tc 行为", "为什么 70K 表现与 60K 和 77K 都不同?")

phys_lines = [
    ("两流体模型 (Two-Fluid Model)", True, Pt(16), ACCENT),
    ("", False, Pt(6), WHITE),
    ("超导电子密度:  ns(T) = ns(0) * [1 - (T/Tc)^4]", False, Pt(14), WHITE),
    ("动能电感:      Lk(T) = Lk(0) / [1 - (T/Tc)^4]", False, Pt(14), WHITE),
    ("谐振频率:      f0(T) = f0(0) / sqrt(1 + Lk(T)/Lg)", False, Pt(14), WHITE),
    ("品质因子:      Qi(T) ∝ ns(T) ∝ [1 - (T/Tc)^4]", False, Pt(14), WHITE),
    ("", False, Pt(8), WHITE),
    ("当 T << Tc (6-60K):", True, Pt(14), GREEN),
    ("  ns(T) 变化平缓 → df/dT 小且恒定 (~-15 MHz/K), dip depth 高且稳定", False, Pt(12), WHITE),
    ("  → 60K 以下所有谐振子行为一致, 数据质量 A 级", False, Pt(12), WHITE),
    ("", False, Pt(6), WHITE),
    ("当 T 接近 Tc (70-77K):", True, Pt(14), ORANGE),
    ("  ns(T) 加速下降 → df/dT 翻倍至 ~25-33 MHz/K, Qi 显著降低", False, Pt(12), WHITE),
    ("  ns(T)/ns(0) 在 77K 仅约 0.35, 在 80K 约 0.20", False, Pt(12), WHITE),
    ("  → 所有五个谐振子同步进入 ns 加速衰减区, 行为高度一致", False, Pt(12), WHITE),
    ("", False, Pt(6), WHITE),
    ("70K 修正后的关键发现:", True, Pt(14), ACCENT2),
    ("  修正 R1/R2 识别互换后, 五个谐振子 df/dT 翻倍率高度一致 (1.96-2.07x)", False, Pt(12), WHITE),
    ("  这证明: 近 Tc 行为是整体均匀的, 并非个别谐振子的异常", False, Pt(12), WHITE),
    ("  R5 的 dip 暴跌 (26→4.4 dB) 是唯一突出的个体差异", False, Pt(12), WHITE),
    ("  → R5 所在区域有效 Tc 可能最低, 或其电流密度分布对耗散最敏感", False, Pt(12), WHITE),
    ("", False, Pt(6), WHITE),
    ("结论: 修正追踪错误后, 70K 数据高度物理自洽", True, Pt(13), ACCENT),
    ("算法在 70K 互换了 R1/R2 识别, 修正后迹象消失——这提醒我们,", False, Pt(12), LGRAY),
    ("高温段必须人工验证追踪结果, 不可盲信自动算法。", False, Pt(12), LGRAY),
]
_ml(s, Inches(0.8), Inches(2.2), Inches(11.5), Inches(5.2), phys_lines)
_pn(s, 7)

# ============================================================
# SLIDE 8: R5 穿越寄生坑 — S21 轨迹图
# ============================================================
s = _add_blank_slide()
_add_bg(s, DARK)
_add_dark_header(s, "R5 穿越固定频率寄生坑", "所有温度 S21 叠加: R5 从 5.25 GHz 红移至 4.66 GHz, 途中穿过 ~5.09 GHz 寄生模式")

_img(s, IMG_DIR / "R5_trajectory_vs_parasitic.png",
     Inches(0.3), Inches(2.1), Inches(12.5))

_ml(s, Inches(0.8), Inches(6.7), Inches(11.5), Inches(0.6), [
    ("Red dashed line: fixed-frequency parasitic at ~5.09 GHz (present at ALL temperatures, does NOT shift with T). Black circles: R5 trajectory.",
     False, Pt(10), LGRAY),
])
_pn(s, 8)

# ============================================================
# SLIDE 9: 寄生坑分析 — Prominence 交叉
# ============================================================
s = _add_blank_slide()
_add_bg(s, DARK)
_add_dark_header(s, "R5 vs 寄生坑: Prominence 交叉与误追踪风险", "50-60K 时 R5 与寄生间距 < 50 MHz, 追踪器无法分辨")

# Left: prominence comparison
_img(s, IMG_DIR / "R5_vs_parasitic_prominence.png",
     Inches(0.3), Inches(2.1), Inches(7.3))

# Right: detail image
_img(s, IMG_DIR / "parasitic_dip_detail.png",
     Inches(7.8), Inches(2.1), Inches(5.2))

# Commentary
_ml(s, Inches(7.8), Inches(5.6), Inches(5.2), Inches(1.5), [
    ("关键发现:", True, Pt(13), ORANGE),
    ("R5 的 dip 暴跌 (26→4.4 dB) 并非物理退化, 而是穿过了寄生坑。", False, Pt(11), WHITE),
    ("50-60K 时 R5 与寄生重叠, prominence 虚高 (26-33 dB 是叠加值)。", False, Pt(11), WHITE),
    ("70K 后 R5 已穿过寄生, prominence 骤降至 2.1 dB 才是真实值。", False, Pt(11), WHITE),
    ("寄生在 5.086-5.098 GHz 区间, 完全不随温度移动 → 盒子模式/电缆谐振。", False, Pt(11), WHITE),
    ("R5 的 δf/f₀ 在 50-70K 段因混淆寄生, 可靠性有限。", False, Pt(11), LGRAY),
])
_pn(s, 9)

# ============================================================
# SLIDE 10: Dip Depth Comparison
# ============================================================
s = _add_blank_slide()
_add_bg(s, DARK)
_add_dark_header(s, "谐振深度对比: OLD vs NEW", "Dip depth 一致性表明器件品质因子未退化")

_img(s, IMG_DIR / "dip_depth_comparison.png",
     Inches(0.5), Inches(2.1), Inches(12.3))

# Note
_ml(s, Inches(0.8), Inches(6.6), Inches(11.5), Inches(0.6), [
    ("Note: R5 excluded from dip depth comparison due to parasitic contamination at 50-70K. Other 4 resonators show excellent agreement.",
     False, Pt(11), LGRAY),
])
_pn(s, 10)

# ============================================================
# SLIDE 11: Key Findings
# ============================================================
s = _add_blank_slide()
_add_bg(s, DARK)
_add_dark_header(s, "主要发现与结论")

findings_lines = [
    ("1. 器件稳定性确认", True, Pt(16), ACCENT),
    ("  两次测量 (间隔 ~2 周) 的谐振 dip depth 几乎完全一致，证明 YBCO KID 器件未发生退化。", False, Pt(13), WHITE),
    ("", False, Pt(8), WHITE),
    ("2. 低温段 (6-20K) 完美吻合", True, Pt(16), GREEN),
    ("  f0 差异 < 0.3 MHz，测量可重复性极好。", False, Pt(13), WHITE),
    ("", False, Pt(8), WHITE),
    ("3. 高温段 (40-70K) 存在系统性 f0 偏移", True, Pt(16), ORANGE),
    ("  OLD 测量频率系统偏高 3-25 MHz → OLD 样品温度偏低约 2-10K。", False, Pt(13), WHITE),
    ("  可能原因: LakeShore 335 温度计校准漂移或热锚差异。", False, Pt(13), WHITE),
    ("", False, Pt(8), WHITE),
    ("4. R5 穿越固定频率寄生坑 (重要发现!)", True, Pt(16), ACCENT2),
    ("  在 ~5.09 GHz 存在一个完全不随温度移动的深 dip (盒子模式/电缆谐振)。", False, Pt(13), WHITE),
    ("  R5 从 5.25 GHz 红移至 4.66 GHz, 在 50-60K 之间横穿此寄生坑。", False, Pt(13), WHITE),
    ("  R5 的 dip 暴跌 (26→4.4 dB) 是穿越寄生效应的结果, 之前的高 prominence 是叠加假象。", False, Pt(13), WHITE),
    ("  R5 在 50-70K 段的 delta_f/f0 数据受寄生污染, 需谨慎使用。", False, Pt(13), WHITE),
    ("", False, Pt(8), WHITE),
    ("5. 70K R1/R2 算法互换已修正", True, Pt(16), ACCENT),
    ("  修正后五个谐振子 df/dT 翻倍率高度一致 (1.96-2.07x) → 近 Tc 行为整体均匀。", False, Pt(13), WHITE),
    ("", False, Pt(8), WHITE),
    ("6. NEW 数据集的优势", True, Pt(16), ACCENT),
    ("  16 级 VNA 功率 (2 dB 步长) vs OLD 的 3 级 — 可做精细的功率依赖分析。", False, Pt(13), WHITE),
]
_ml(s, Inches(0.8), Inches(2.2), Inches(11.5), Inches(5.2), findings_lines)
_ml(s, Inches(0.8), Inches(2.2), Inches(11.5), Inches(5.2), findings_lines)
_pn(s, 11)

# ============================================================
# SLIDE 12: 附录 — 数据目录索引
# ============================================================
s = _add_blank_slide()
_add_bg(s, DARK)
_add_dark_header(s, "附录: 数据与图表索引")

appendix_lines = [
    ("合并数据集", True, Pt(14), ACCENT),
    ("  experiment_data/~merged/20260609-0624__6-80K__full/", False, Pt(11), WHITE),
    ("", False, Pt(4), WHITE),
    ("旧数据集 (对比基准)", True, Pt(14), ACCENT),
    ("  experiment_data/accomplish_merged/", False, Pt(11), WHITE),
    ("", False, Pt(4), WHITE),
    ("验证图", True, Pt(14), ACCENT),
    ("  ~merged/output/_cache/verification_20260609-0624__6-80K__full/", False, Pt(11), WHITE),
    ("", False, Pt(4), WHITE),
    ("画图输出 (方案 A/B/S21)", True, Pt(14), ACCENT),
    ("  ~merged/output/_cache/plot_output/", False, Pt(11), WHITE),
    ("", False, Pt(4), WHITE),
    ("对比分析图", True, Pt(14), ACCENT),
    ("  Data_process/output/comparison/", False, Pt(11), WHITE),
    ("", False, Pt(4), WHITE),
    ("Junk 目录 (17 个空壳/失败运行)", True, Pt(14), ACCENT),
    ("  experiment_data/_junk/", False, Pt(11), WHITE),
]
_ml(s, Inches(0.8), Inches(2.2), Inches(11.5), Inches(5.0), appendix_lines)
_pn(s, 12)

# ============================================================
# Save
# ============================================================
prs.save(str(PPTX_OUT))
print(f"[OK] PPT saved to: {PPTX_OUT}")
print(f"  Slides: {len(prs.slides)}")
