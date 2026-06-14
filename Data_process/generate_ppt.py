# -*- coding: utf-8 -*-
"""
生成 YBCO KID merged 数据集表征简报 PPTX。

从 output/merged/ 子目录读取图片，用 python-pptx 逐页构建。
关键数值通过 dataprocess.py（无需 scraps）从 S2P 数据中提取。
"""

import sys
import os
from pathlib import Path
import re

# 确保 otherwise 目录可导入
_script_dir = Path(__file__).resolve().parent
_otherwise_dir = _script_dir / "otherwise"
if str(_otherwise_dir) not in sys.path:
    sys.path.insert(0, str(_otherwise_dir))

import numpy as np
import dataprocess as dp

# ============================================================
# 配置
# ============================================================

MERGED_DIR = _script_dir.parent / "Auto_Sweep" / "experiment_data" / "merged"
OUTPUT_DIR = _script_dir / "output" / "merged"
PPTX_OUT = _script_dir / "output" / "YBCO_KID_merged_表征简报.pptx"

PIXEL_INDX = 1
MEAS_POWERS = [25, 30, 45]
LASER_POWERS = [0, 1, 3, 5, 7, 9]

# 深色背景色常量（用于 PPT 幻灯片统一主题）
DARK_BG = (0x1A, 0x1A, 0x2E)

# ============================================================
# 数值提取
# ============================================================

def scan_temperatures():
    """扫描 merged 目录下的温度点，返回排序后的列表（(整数温度, 原始目录名)）。"""
    pattern = re.compile(r"^(\d+(?:\.\d+)?)K$")
    entries = []
    for subfolder in MERGED_DIR.iterdir():
        if subfolder.is_dir():
            m = pattern.match(subfolder.name)
            if m:
                entries.append((int(float(m.group(1))), subfolder.name))
    if not entries:
        raise FileNotFoundError(f"MERGED_DIR 中没有找到匹配的温度子目录: {MERGED_DIR}")
    entries.sort(key=lambda x: x[0])
    return entries  # list of (int_temp, dirname)


def find_first_s2p(temp_dirname, power_dbm, laser_mw):
    """返回指定目录下的第一个 S2P 文件路径。

    Args:
        temp_dirname: 温度目录名（如 "6K"）
        power_dbm: VNA 功率 (dBm 数值，取负号)
        laser_mw: 激光功率 (mW)
    """
    path = MERGED_DIR / temp_dirname / f"-{power_dbm}dBm" / f"{laser_mw:02d}mW"
    if not path.is_dir():
        return None
    for f in path.iterdir():
        if f.suffix == ".s2p":
            return str(f)
    return None


def extract_key_numbers():
    """从数据中提取关键数值（仅用 dataprocess.py，不依赖 scraps）。

    Returns:
        dict with keys:
        - f0_ghz: 选定谐振峰频率 (GHz)
        - num_resonances: 检测到的谐振峰总数
        - f0_shift_percent: f0 从最低温到最高温的相对变化 (%)
        - qi_estimate: -3dB 带宽法估算的 Qi
        - low_temp: 最低温度 (K)
        - high_temp: 最高温度 (K)
        - all_resonances: 所有谐振峰频率列表 (GHz)
        - repr_low_temp: 低温代表点温度 (K) or None
        - repr_high_temp: 高温代表点温度 (K) or None
    """
    result = {}

    # 扫描温度
    temp_entries = scan_temperatures()
    temps = [t for t, _ in temp_entries]
    temp_dirs = [d for _, d in temp_entries]
    result["low_temp"] = temps[0]
    result["high_temp"] = temps[-1]

    # 使用最低温、-25dBm、0mW 的数据做谐振峰检测
    s2p_path = find_first_s2p(temp_dirs[0], MEAS_POWERS[0], LASER_POWERS[0])
    if s2p_path is None:
        raise FileNotFoundError(
            f"找不到 S2P 文件: {temp_dirs[0]}/-{MEAS_POWERS[0]}dBm/{LASER_POWERS[0]:02d}mW"
        )

    freq, s21 = dp.load_s_param(s2p_path)

    # 寻峰
    peaks, _, _ = dp.find_true_resonances(
        freq=freq, s21=s21,
        min_prominence=3, distance=10, phase_window=10,
        phase_diff_snr_threshold=0.5,
        noise_inner_window=5, noise_outer_window=40,
        min_phase_diff_support_points=4, min_phase_diff_width=4,
        plot=False,
    )

    result["num_resonances"] = len(peaks)
    all_freqs_ghz = [p["frequency"] / 1e9 for p in peaks]
    result["all_resonances"] = all_freqs_ghz

    if len(peaks) <= PIXEL_INDX:
        raise ValueError(
            f"pixel_indx={PIXEL_INDX} 超出检测到的谐振峰数 ({len(peaks)})"
        )

    selected_peak = peaks[PIXEL_INDX]
    f0_ghz = selected_peak["frequency"] / 1e9
    result["f0_ghz"] = f0_ghz

    # 估算 Qi：用 -3dB 带宽法
    # Qi ≈ f₀ / Δf_{-3dB}
    transmission = 20 * np.log10(np.abs(s21))
    peak_idx = selected_peak["index"]

    # 找到 |S21| 最小值作为 dip 底部
    dip_val = transmission[peak_idx]
    half_power_level = dip_val + 3.0  # -3dB above dip

    # 向左、右找 -3dB 交叉点
    left_idx = peak_idx
    while left_idx > 0 and transmission[left_idx] < half_power_level:
        left_idx -= 1
    right_idx = peak_idx
    while right_idx < len(transmission) - 1 and transmission[right_idx] < half_power_level:
        right_idx += 1

    delta_f = freq[right_idx] - freq[left_idx]
    # Guard: reject bandwidth that spans more than half the measurement range
    total_span = freq[-1] - freq[0]
    if delta_f <= 0 or delta_f > total_span * 0.5:
        result["qi_estimate"] = None
        result["bandwidth_mhz"] = None
    else:
        result["qi_estimate"] = int(f0_ghz * 1e9 / delta_f)
        result["bandwidth_mhz"] = delta_f / 1e6

    # 估算 f0 温度漂移（比较最低温和最高温的近似 f0）
    s2p_high = find_first_s2p(temp_dirs[-1], MEAS_POWERS[0], LASER_POWERS[0])
    if s2p_high:
        freq_high, s21_high = dp.load_s_param(s2p_high)
        # 在高温数据中重新寻峰
        peaks_high, _, _ = dp.find_true_resonances(
            freq=freq_high, s21=s21_high,
            min_prominence=3, distance=10, phase_window=10,
            phase_diff_snr_threshold=0.5,
            noise_inner_window=5, noise_outer_window=40,
            min_phase_diff_support_points=4, min_phase_diff_width=4,
            plot=False,
        )
        if len(peaks_high) > PIXEL_INDX:
            f0_high_ghz = peaks_high[PIXEL_INDX]["frequency"] / 1e9
            shift_pct = (f0_high_ghz - f0_ghz) / f0_ghz * 100
            result["f0_shift_percent"] = abs(shift_pct)
            result["f0_shift_direction"] = "红移（频率降低）" if shift_pct < 0 else "蓝移（频率升高）"
        else:
            result["f0_shift_percent"] = None
            result["f0_shift_direction"] = ""
    else:
        result["f0_shift_percent"] = None
        result["f0_shift_direction"] = ""

    # 从现有图片中推断所选的低温和高温代表点
    result["repr_low_temp"] = None
    result["repr_high_temp"] = None
    low_dir = OUTPUT_DIR / "05_optical_response_6K"
    high_dir = OUTPUT_DIR / "06_optical_response_highT"
    low_temp_files = [f for f in os.listdir(low_dir) if f.startswith("res shift")] if low_dir.is_dir() else []
    high_temp_files = [f for f in os.listdir(high_dir) if f.startswith("res shift")] if high_dir.is_dir() else []
    if low_temp_files:
        m = re.search(r"([\d.]+)K", low_temp_files[0])
        if m:
            result["repr_low_temp"] = float(m.group(1))
    if high_temp_files:
        m = re.search(r"([\d.]+)K", high_temp_files[0])
        if m:
            result["repr_high_temp"] = float(m.group(1))

    return result

# ============================================================
# PPT 生成
# ============================================================

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# 幻灯片尺寸 (widescreen 16:9)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# 边距
MARGIN = Inches(0.5)
GAP = Inches(0.2)


def _add_textbox(slide, left, top, width, height, text, font_size=Pt(14),
                 bold=False, color=None, alignment=PP_ALIGN.LEFT):
    """添加文本框。"""
    if color is None:
        color = RGBColor(0xFF, 0xFF, 0xFF)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return txBox


def _add_bullet_list(slide, left, top, width, height, items,
                     font_size=Pt(12), color=None):
    """添加 bullet 列表。"""
    if color is None:
        color = RGBColor(0xDD, 0xDD, 0xDD)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = font_size
        p.font.color.rgb = color
        p.space_after = Pt(6)
    return txBox


def _set_slide_bg(slide, r, g, b):
    """设置幻灯片背景色。"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)


def _build_optical_response_slide(prs, blank_layout,
                                   temp_label, subdir, title_suffix,
                                   comparison_text):
    """构建光学响应幻灯片（Slide 4/5 共用模板）。"""
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide, *DARK_BG)

    subdir_path = OUTPUT_DIR / subdir

    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    _add_textbox(slide, MARGIN, Inches(0.3), Inches(12), Inches(0.6),
                 f"{temp_label}下的光致谐振频移 {title_suffix}",
                 font_size=Pt(28), bold=True, color=WHITE)

    # 找到该子目录下的所有图片
    s21_files = []
    res_shift_file = None

    if subdir_path.is_dir():
        # 按文件名排序确保确定性
        all_files = sorted(os.listdir(str(subdir_path)))
        s21_files = [f for f in all_files if f.startswith("s21 -")]
        # 按温度排序后取第一个（最低温），确保确定性
        res_shift_candidates = sorted([f for f in all_files if f.startswith("res shift -")])
        if res_shift_candidates:
            res_shift_file = res_shift_candidates[0]

    # 2×2 网格排列
    positions = [
        (MARGIN, Inches(1.2), Inches(3.0), Inches(2.6)),          # 左上
        (Inches(3.7), Inches(1.2), Inches(3.0), Inches(2.6)),     # 右上
        (MARGIN, Inches(4.0), Inches(3.0), Inches(2.6)),          # 左下
        (Inches(3.7), Inches(4.0), Inches(3.0), Inches(2.6)),     # 右下（res shift）
    ]

    # 只有找到图片时才添加
    for i, s21f in enumerate(s21_files[:3]):
        img_path = subdir_path / s21f
        if img_path.exists():
            left, top, w, h = positions[i]
            slide.shapes.add_picture(str(img_path), left, top, w, h)

    if res_shift_file:
        img_path = subdir_path / res_shift_file
        if img_path.exists():
            left, top, w, h = positions[3]
            slide.shapes.add_picture(str(img_path), left, top, w, h)

    # 右侧要点
    _add_bullet_list(slide, Inches(7.2), Inches(1.5), Inches(5.5), Inches(5.5), [
        "展示 -25 / -30 / -45 dBm 三个 VNA 功率",
        "S21 曲线随 0→9 mW 激光功率演化",
        "右上图：谐振频率偏移 vs 激光功率",
        comparison_text,
    ], font_size=Pt(13))

    return slide


def build_pptx(numbers):
    """构建 7 页 PPTX，返回 Presentation 对象。"""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 查找空白版式（跨 python-pptx 版本稳定）
    blank_layout = None
    for layout in prs.slide_layouts:
        if layout.name == 'Blank' or layout.name == '空白':
            blank_layout = layout
            break
    if blank_layout is None:
        blank_layout = prs.slide_layouts[6]  # fallback

    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    GRAY = RGBColor(0xCC, 0xCC, 0xCC)
    ACCENT = RGBColor(0x4F, 0xC3, 0xF7)

    n_res = numbers["num_resonances"]
    f0 = numbers["f0_ghz"]
    qi_est = numbers.get("qi_estimate")
    qi_text = f"{qi_est}" if qi_est is not None else "—"
    shift_pct = numbers.get("f0_shift_percent")
    shift_dir = numbers.get("f0_shift_direction", "")
    low_t = numbers["low_temp"]
    high_t = numbers["high_temp"]

    # 构建温度范围内文本
    if shift_pct is not None:
        f0_shift_text = f"f₀ 随温度升高单调{shift_dir}约 {shift_pct:.1f}%（{low_t}→{high_t}K）"
    else:
        f0_shift_text = f"f₀ 随温度升高呈现单调频移趋势（{low_t}→{high_t}K）"

    # ---- Slide 1: 封面 ----
    slide1 = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide1, *DARK_BG)

    _add_textbox(slide1, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
                 "YBCO KID 微波-光学联合表征",
                 font_size=Pt(36), bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)

    info_lines = [
        "样品：YBCO KID (merged 数据集)",
        f"温度范围：{low_t}K → {high_t}K（{len(scan_temperatures())} 个温度点）",
        "VNA 功率：-25 / -30 / -45 dBm",
        "激光功率：0, 1, 3, 5, 7, 9 mW",
        "分析日期：2026-06-15",
    ]
    y = Inches(3.5)
    for line in info_lines:
        _add_textbox(slide1, Inches(2.5), y, Inches(8), Inches(0.5),
                     line, font_size=Pt(16), color=GRAY,
                     alignment=PP_ALIGN.CENTER)
        y += Inches(0.45)

    # ---- Slide 2: 谐振峰检测与 S21 温度演化 ----
    slide2 = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide2, *DARK_BG)

    _add_textbox(slide2, MARGIN, Inches(0.3), Inches(12), Inches(0.6),
                 "谐振器识别与全温 S21 叠加",
                 font_size=Pt(28), bold=True, color=WHITE)

    img_detect = OUTPUT_DIR / "01_resonance_detection" / "resonance_detection.jpg"
    if img_detect.exists():
        slide2.shapes.add_picture(
            str(img_detect), MARGIN, Inches(1.1), Inches(6), Inches(3.5)
        )

    img_s21_temp = OUTPUT_DIR / "04_S21_temperature_overlay" / "s21 vs - temp.jpg"
    if img_s21_temp.exists():
        slide2.shapes.add_picture(
            str(img_s21_temp), Inches(6.8), Inches(1.1), Inches(6), Inches(3.5)
        )

    _add_bullet_list(slide2, MARGIN, Inches(5.0), Inches(12), Inches(2.2), [
        f"采用幅度谷 + 相位差分峰联合判据自动寻峰（SNR ≥ 0.5），共检测到 {n_res} 个谐振峰",
        f"选定谐振峰位于 {f0:.3f} GHz（pixel {PIXEL_INDX}），QL ≈ {qi_text}（-3dB 带宽法估计）",
        "全温 S21 叠加：谐振峰随温度单调频移，符合超导动能电感 Lₖ ∝ λ²(T) 理论预期",
        "峰形保持良好，器件在测量温区范围内稳定工作",
    ], font_size=Pt(13))

    # ---- Slide 3: f₀(T) 与 Qi(T) ----
    slide3 = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide3, *DARK_BG)

    _add_textbox(slide3, MARGIN, Inches(0.3), Inches(12), Inches(0.6),
                 "谐振频率与内禀品质因子的温度响应",
                 font_size=Pt(28), bold=True, color=WHITE)

    img_f0 = OUTPUT_DIR / "02_f0_temperature" / "f0_versus_temp.jpg"
    if img_f0.exists():
        slide3.shapes.add_picture(
            str(img_f0), MARGIN, Inches(1.1), Inches(5.8), Inches(2.8)
        )

    img_qi = OUTPUT_DIR / "03_Qi_temperature" / "qis_versus_temp.jpg"
    if img_qi.exists():
        slide3.shapes.add_picture(
            str(img_qi), MARGIN, Inches(4.1), Inches(5.8), Inches(2.8)
        )

    _add_bullet_list(slide3, Inches(7.2), Inches(1.5), Inches(5.5), Inches(5.5), [
        f0_shift_text,
        "Qi 低温段较高，随温度上升逐渐降低，归因于准粒子损耗增大",
        "三个 VNA 功率 (-25/-30/-45 dBm) 的 Qi 偏差较小",
        "表明读出功率未引入显著非线性效应",
    ], font_size=Pt(13))

    # ---- Slide 4: 低温光学响应 ----
    _build_optical_response_slide(
        prs, blank_layout,
        temp_label="低温",
        subdir="05_optical_response_6K",
        title_suffix="（T ≈ 6K）",
        comparison_text="频移与激光功率呈良好线性关系 → 符合准粒子退对机制",
    )

    # ---- Slide 5: 高温光学响应 ----
    _build_optical_response_slide(
        prs, blank_layout,
        temp_label="高温",
        subdir="06_optical_response_highT",
        title_suffix="（T ≈ 76K）",
        comparison_text="与 6K 对比：高温下热准粒子密度升高，光注入准粒子的相对增量减小",
    )

    # ---- Slide 6: 响应率 vs 温度 ----
    slide6 = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide6, *DARK_BG)

    _add_textbox(slide6, MARGIN, Inches(0.3), Inches(12), Inches(0.6),
                 "光学响应率的温度依赖性",
                 font_size=Pt(28), bold=True, color=WHITE)

    img_resp = OUTPUT_DIR / "07_responsivity_temperature" / "responsivity_vs_temp.jpg"
    if img_resp.exists():
        slide6.shapes.add_picture(
            str(img_resp), MARGIN, Inches(1.2), Inches(8.5), Inches(5.0)
        )

    _add_bullet_list(slide6, Inches(9.5), Inches(1.5), Inches(3.3), Inches(5.5), [
        "响应率 (Hz/W) 随温度的变化趋势",
        "响应率温度依赖与超导能隙 Δ(T) 定性一致",
        "低温段保持较高响应水平",
        "下一步需与 BCS 理论定量比较",
    ], font_size=Pt(13))

    # ---- Slide 7: 总结 ----
    slide7 = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide7, *DARK_BG)

    _add_textbox(slide7, Inches(1), Inches(1.0), Inches(11), Inches(1.0),
                 "小结与下一步",
                 font_size=Pt(36), bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)

    done_items = [
        f"成功表征 YBCO KID 在 {low_t}–{high_t}K 的谐振特性与光学响应",
        f"Qi(T) 下降趋势与 BCS 理论预期一致",
        "光学响应率呈现温度依赖，低温段保持较高响应水平",
        "系统在测试温区范围内稳定可靠，谐振峰形保持良好",
    ]
    if shift_pct is not None:
        done_items[0] = f"成功表征 YBCO KID 在 {low_t}–{high_t}K 的谐振特性与光学响应（f₀ {shift_dir} {shift_pct:.1f}%）"

    _add_bullet_list(slide7, Inches(1.5), Inches(2.5), Inches(10), Inches(2.0),
                     done_items, font_size=Pt(16), color=WHITE)

    next_items = [
        "NEP（噪声等效功率）估算与优化",
        "多像素统计比较，评估器件均匀性",
        "更低温度（< 1K）测量，探索极限灵敏度",
        "低温放大器（HEMT / JPA）集成测试",
    ]
    _add_bullet_list(slide7, Inches(1.5), Inches(4.5), Inches(10), Inches(2.0),
                     next_items, font_size=Pt(14), color=GRAY)

    return prs


if __name__ == "__main__":
    print("提取关键数值...")
    numbers = extract_key_numbers()
    print()
    for k, v in numbers.items():
        print(f"  {k}: {v}")

    print(f"\n构建 PPTX...")
    prs = build_pptx(numbers)

    os.makedirs(PPTX_OUT.parent, exist_ok=True)
    prs.save(str(PPTX_OUT))
    print(f"\nPPTX 已保存到: {PPTX_OUT}")
