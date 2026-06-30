# -*- coding: utf-8 -*-
"""
生成 YBCO KID 五谐振器对比表征简报 v6。

使用 output/merged/compare/ 下 PPT 优化后的图片，白底 16:9 版式。
"""

import csv
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# ============================================================
# 路径配置
# ============================================================
SCRIPT_DIR = Path(__file__).resolve().parent
IMG_DIR = SCRIPT_DIR / "output" / "merged" / "compare"
CSV_PATH = IMG_DIR / "resonator_summary.csv"
PPTX_OUT = SCRIPT_DIR / "output" / "YBCO_KID_五谐振器对比_v6.pptx"

# ============================================================
# 颜色 / 字体
# ============================================================
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK  = RGBColor(0x1A, 0x1A, 0x1A)
GRAY  = RGBColor(0x66, 0x66, 0x66)
LGRAY = RGBColor(0xAA, 0xAA, 0xAA)
ACCENT = RGBColor(0x1F, 0x77, 0xB4)  # Tableau blue, matches plot colors
ACCENTS = [
    RGBColor(0x1F, 0x77, 0xB4),  # blue
    RGBColor(0xD6, 0x27, 0x28),  # red
    RGBColor(0x2C, 0xA0, 0x2C),  # green
    RGBColor(0xFF, 0x7F, 0x0E),  # orange
    RGBColor(0x94, 0x67, 0xBD),  # purple
]
LIGHT_BG = RGBColor(0xF5, 0xF5, 0xF5)

FONT_TITLE = "Arial"
FONT_BODY  = "Arial"

# ============================================================
# 辅助函数
# ============================================================

def load_summary():
    """读取 resonator_summary.csv 返回字典列表。"""
    rows = []
    with open(CSV_PATH, "r", encoding="utf-8") as f:
        reader = csv.DictReader(f)
        for r in reader:
            rows.append(r)
    return rows


def _add_textbox(slide, left, top, width, height, text,
                 font_size=Pt(14), bold=False, color=DARK,
                 alignment=PP_ALIGN.LEFT, font_name=FONT_BODY,
                 anchor=MSO_ANCHOR.TOP):
    """添加文本框。"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    txBox.text_frame.auto_size = None
    tf = txBox.text_frame
    tf.paragraphs[0].alignment = alignment
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    return txBox


def _add_multiline(slide, left, top, width, height, lines,
                   font_size=Pt(12), color=DARK, line_spacing=Pt(20),
                   font_name=FONT_BODY):
    """添加多行文本框，每行一个 paragraph。"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, bold, fs, c) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = fs if fs else font_size
        p.font.bold = bold
        p.font.color.rgb = c if c else color
        p.font.name = font_name
        p.space_after = line_spacing
    return txBox


def _add_image_safe(slide, img_name, left, top, width, height):
    """安全添加图片，文件不存在时返回 None。"""
    path = IMG_DIR / img_name
    if path.exists():
        return slide.shapes.add_picture(str(path), left, top, width, height)
    print(f"  [WARN] 图片不存在: {path}")
    return None


def _set_slide_bg(slide, r, g, b):
    """设置幻灯片纯色背景。"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)


def _add_page_number(slide, num, total):
    """右下角页码。"""
    _add_textbox(slide, Inches(12.0), Inches(7.0), Inches(1.0), Inches(0.3),
                 f"{num} / {total}", font_size=Pt(9), color=LGRAY,
                 alignment=PP_ALIGN.RIGHT)


def _add_accent_line(slide, left, top, width):
    """添加装饰性细线。"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Pt(2.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    return shape


# ============================================================
# 幻灯片构建
# ============================================================

def build_pptx():
    summary = load_summary()

    # ---- Presentation setup ----
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]  # blank

    TOTAL = 12
    page = [0]  # mutable counter

    def new_slide():
        page[0] += 1
        slide = prs.slides.add_slide(blank)
        _set_slide_bg(slide, 0xFF, 0xFF, 0xFF)
        return slide

    # ================================================================
    # Slide 1: 封面
    # ================================================================
    s = new_slide()
    _set_slide_bg(s, 0xF8, 0xF8, 0xF8)

    # 顶部装饰条
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.08))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()

    _add_textbox(s, Inches(1.5), Inches(1.5), Inches(10.3), Inches(1.2),
                 "YBCO 片上五谐振器\n微波-光学联合表征",
                 font_size=Pt(40), bold=True, color=DARK,
                 alignment=PP_ALIGN.CENTER, font_name=FONT_TITLE)

    _add_accent_line(s, Inches(5.5), Inches(3.2), Inches(2.3))

    # 副标题
    info = [
        (f"赵思源", True, Pt(20), DARK),
        ("", False, Pt(8), DARK),
        (f"频率覆盖: 3.85 – 5.25 GHz    温度范围: 6 K → 80 K    VNA: −25 dBm    激光: 0–9 mW", False, Pt(14), GRAY),
        ("", False, Pt(6), DARK),
        (f"5 枚 hanger 型 CPW 谐振器 · YBCO 高温超导薄膜 · 2026-06-18", False, Pt(12), LGRAY),
    ]
    _add_multiline(s, Inches(2.0), Inches(3.6), Inches(9.3), Inches(2.0),
                   info, line_spacing=Pt(14))

    _add_page_number(s, page[0], TOTAL)

    # ================================================================
    # Slide 2: 全谱指纹
    # ================================================================
    s = new_slide()
    _add_textbox(s, Inches(0.6), Inches(0.3), Inches(12.0), Inches(0.6),
                 "全谱指纹：五谐振器 S21 全景", font_size=Pt(28), bold=True, color=DARK)
    _add_accent_line(s, Inches(0.6), Inches(0.85), Inches(1.5))

    _add_image_safe(s, "01_spectrum_overview.jpg",
                    Inches(0.6), Inches(1.1), Inches(9.5), Inches(4.5))

    # 右侧参数表
    cmts = [
        ("参数  @ 6 K, −25 dBm, 0 mW", True, Pt(16), DARK),
        ("", False, Pt(6), DARK),
        ("谐振器       f₀ / GHz      Dip / dB       QL", True, Pt(12), GRAY),
    ]
    for i, row in enumerate(summary):
        c = ACCENTS[i % len(ACCENTS)]
        cmts.append((f"R{i+1}             {row['f0_6K_GHz']}          {row['dip_6K_dB']}           {row['QL_6K']}",
                     False, Pt(13), c))
    cmts.append(("", False, Pt(8), DARK))
    cmts.append(("R4 最深耦合 (−17.7 dB), R3 最高 QL (4167)", False, Pt(11), GRAY))
    cmts.append(("5 个谐振器均匀分布在 3.85–5.25 GHz 馈线上", False, Pt(11), GRAY))

    _add_multiline(s, Inches(10.5), Inches(1.1), Inches(2.5), Inches(5.5),
                   cmts, line_spacing=Pt(10))
    _add_page_number(s, page[0], TOTAL)

    # ================================================================
    # Slide 3: f₀(T) 五线对比
    # ================================================================
    s = new_slide()
    _add_textbox(s, Inches(0.6), Inches(0.3), Inches(12.0), Inches(0.6),
                 "谐振频率温度响应：f₀(T) 五线对比", font_size=Pt(28), bold=True, color=DARK)
    _add_accent_line(s, Inches(0.6), Inches(0.85), Inches(1.5))

    _add_image_safe(s, "02_f0_vs_temp_all.jpg",
                    Inches(0.6), Inches(1.1), Inches(12.0), Inches(5.2))

    # 底部要点
    _add_multiline(s, Inches(0.6), Inches(6.5), Inches(12.0), Inches(0.8), [
        ("所有谐振器随温度单调红移（频率降低），频移幅度一致 → YBCO 薄膜均匀，Tc 一致。动能电感 Lₖ ∝ λ²(T) = λ₀/√(1−(T/Tc)⁴) 驱动频移。R1–R4 可追踪至 ~70 K，R5 (5.252 GHz) 近带边，仅至 ~28 K。",
         False, Pt(12), GRAY),
    ], line_spacing=Pt(4))
    _add_page_number(s, page[0], TOTAL)

    # ================================================================
    # Slide 4: QL(T) 五线对比
    # ================================================================
    s = new_slide()
    _add_textbox(s, Inches(0.6), Inches(0.3), Inches(12.0), Inches(0.6),
                 "品质因子温度响应：QL(T) 五线对比", font_size=Pt(28), bold=True, color=DARK)
    _add_accent_line(s, Inches(0.6), Inches(0.85), Inches(1.5))

    _add_image_safe(s, "03_qi_vs_temp_all.jpg",
                    Inches(0.6), Inches(1.1), Inches(12.0), Inches(5.2))

    _add_multiline(s, Inches(0.6), Inches(6.5), Inches(12.0), Inches(0.8), [
        ("R3 (4.500 GHz) QL 最高 ≈ 4167；R4 (4.997 GHz) 耦合最深 (−17.7 dB)，QL 仍保持 ≈ 3966。QL 随温度上升单调下降 → 热准粒子密度增加 → 损耗增大。低温段 QL 值接近 → 剩余损耗由耦合主导而非材料。",
         False, Pt(12), GRAY),
    ], line_spacing=Pt(4))
    _add_page_number(s, page[0], TOTAL)

    # ================================================================
    # Slide 5: 6K 光学响应
    # ================================================================
    s = new_slide()
    _add_textbox(s, Inches(0.6), Inches(0.3), Inches(12.0), Inches(0.6),
                 "6 K 光学响应：δf/f₀ vs 激光功率", font_size=Pt(28), bold=True, color=DARK)
    _add_accent_line(s, Inches(0.6), Inches(0.85), Inches(1.5))

    _add_image_safe(s, "04_optical_response_6k_all.jpg",
                    Inches(0.6), Inches(1.1), Inches(12.0), Inches(5.2))

    _add_multiline(s, Inches(0.6), Inches(6.5), Inches(12.0), Inches(0.8), [
        ("δf/f₀ = (f₀(P) − f₀(0)) / f₀(0)，单位 ppm。频移随激光功率线性增加 → 符合光致准粒子拆对机制。R4 响应最大 (−1260 kHz @ 9 mW)，与耦合深度正相关。Dip 深度随激光功率变化微弱 → 光注入主要影响超流密度而非损耗。",
         False, Pt(12), GRAY),
    ], line_spacing=Pt(4))
    _add_page_number(s, page[0], TOTAL)

    # ================================================================
    # Slide 6: 高温光学响应
    # ================================================================
    s = new_slide()
    _add_textbox(s, Inches(0.6), Inches(0.3), Inches(12.0), Inches(0.6),
                 "光学响应率随温度衰减", font_size=Pt(28), bold=True, color=DARK)
    _add_accent_line(s, Inches(0.6), Inches(0.85), Inches(1.5))

    _add_image_safe(s, "10_optical_composite_hight.jpg",
                    Inches(0.6), Inches(1.1), Inches(12.0), Inches(5.2))

    _add_multiline(s, Inches(0.6), Inches(6.5), Inches(12.0), Inches(0.8), [
        ("左图：Δf/f₀ (9 mW − 0 mW) 随温度升高趋近于零 → 热准粒子淹没光注入信号。右图：R4 在各谐振器可追踪的最高温度下的 S21 局部放大 (±15 MHz)。高通处准粒子本底已显著抬高，谐振谷展宽变浅。",
         False, Pt(12), GRAY),
    ], line_spacing=Pt(4))
    _add_page_number(s, page[0], TOTAL)

    # ================================================================
    # Slide 7: QL vs 读出功率
    # ================================================================
    s = new_slide()
    _add_textbox(s, Inches(0.6), Inches(0.3), Inches(12.0), Inches(0.6),
                 "QL vs VNA 读出功率  @ 6 K, 0 mW", font_size=Pt(28), bold=True, color=DARK)
    _add_accent_line(s, Inches(0.6), Inches(0.85), Inches(1.5))

    _add_image_safe(s, "05_ql_vs_power.jpg",
                    Inches(0.6), Inches(1.1), Inches(8.0), Inches(5.2))

    _add_multiline(s, Inches(9.0), Inches(1.5), Inches(4.0), Inches(4.5), [
        ("读出功率扫描 (−45, −30, −25 dBm)", True, Pt(16), DARK),
        ("", False, Pt(6), DARK),
        ("QL 在不同读出功率下基本一致", False, Pt(13), DARK),
        ("→ 无显著的读出功率非线性效应", False, Pt(13), DARK),
        ("→ 器件在线性响应区内工作", False, Pt(13), DARK),
        ("", False, Pt(8), DARK),
        ("R1–R3: QL 随功率略微变化", False, Pt(12), GRAY),
        ("R4: QL 保持 ~4000 基本不变", False, Pt(12), GRAY),
        ("R5: 数据点不足 (部分功率无谐振)", False, Pt(12), GRAY),
        ("", False, Pt(8), DARK),
        ("建议补充 circle-fit 分离 Qc/Qi", False, Pt(11), LGRAY),
    ], line_spacing=Pt(8))
    _add_page_number(s, page[0], TOTAL)

    # ================================================================
    # Slide 8: 归一化光学响应率 vs T
    # ================================================================
    s = new_slide()
    _add_textbox(s, Inches(0.6), Inches(0.3), Inches(12.0), Inches(0.6),
                 "归一化光学响应率 vs 温度", font_size=Pt(28), bold=True, color=DARK)
    _add_accent_line(s, Inches(0.6), Inches(0.85), Inches(1.5))

    _add_image_safe(s, "06_responsivity_vs_temp.jpg",
                    Inches(0.6), Inches(1.1), Inches(12.0), Inches(5.2))

    _add_multiline(s, Inches(0.6), Inches(6.5), Inches(12.0), Inches(0.8), [
        ("Δf/f₀(9 mW) vs T。各谐振器曲线自然截断于各自超导态消失的温度：R4 可追踪至 ~74 K，R1–R3 至 ~40–58 K，R5 至 ~28 K。响应率为负（光致红移），与准粒子拆对 → σ₂ 减小 → 动能电感增大一致。",
         False, Pt(12), GRAY),
    ], line_spacing=Pt(4))
    _add_page_number(s, page[0], TOTAL)

    # ================================================================
    # Slide 9: Dip 深度 vs T
    # ================================================================
    s = new_slide()
    _add_textbox(s, Inches(0.6), Inches(0.3), Inches(12.0), Inches(0.6),
                 "谐振谷深度 vs 温度", font_size=Pt(28), bold=True, color=DARK)
    _add_accent_line(s, Inches(0.6), Inches(0.85), Inches(1.5))

    _add_image_safe(s, "07_dip_vs_temp.jpg",
                    Inches(0.6), Inches(1.1), Inches(12.0), Inches(5.2))

    _add_multiline(s, Inches(0.6), Inches(6.5), Inches(12.0), Inches(0.8), [
        ("红色虚线 = 检测极限 (−1 dB)。谷深随 T 单调变浅 → 热准粒子增多 → 谐振展宽，损耗增大。R4 最深 (−17.7 dB @ 6 K)，坚持最远 (~74 K)；R5 最浅 (−6.1 dB)，最早消失 (~28 K)。谷深是谐振器可工作温区的直观度量。",
         False, Pt(12), GRAY),
    ], line_spacing=Pt(4))
    _add_page_number(s, page[0], TOTAL)

    # ================================================================
    # Slide 10: 归一化 f₀ 偏移 vs T
    # ================================================================
    s = new_slide()
    _add_textbox(s, Inches(0.6), Inches(0.3), Inches(12.0), Inches(0.6),
                 "归一化谐振频率偏移 vs 温度", font_size=Pt(28), bold=True, color=DARK)
    _add_accent_line(s, Inches(0.6), Inches(0.85), Inches(1.5))

    _add_image_safe(s, "08_normalized_f0_vs_temp.jpg",
                    Inches(0.6), Inches(1.1), Inches(12.0), Inches(5.2))

    _add_multiline(s, Inches(0.6), Inches(6.5), Inches(12.0), Inches(0.8), [
        ("(f₀(T) − f₀(6K)) / f₀(6K)，单位 ppm。五条归一化曲线在重叠温区高度重合 → YBCO 薄膜的动能电感温度响应高度一致，Tc 均匀。不同谐振器绝对 f₀ 由几何决定，但 Δf/f₀(T) 由 BCS 超导能隙 Δ(T) 决定 → 重合说明材料均匀。",
         False, Pt(12), GRAY),
    ], line_spacing=Pt(4))
    _add_page_number(s, page[0], TOTAL)

    # ================================================================
    # Slide 11: 6K 光学响应复合图
    # ================================================================
    s = new_slide()
    _add_textbox(s, Inches(0.6), Inches(0.3), Inches(12.0), Inches(0.6),
                 "6 K 光学响应复合分析：δf/f₀ 散点 + 三功率 S21 局部放大", font_size=Pt(28), bold=True, color=DARK)
    _add_accent_line(s, Inches(0.6), Inches(0.85), Inches(1.5))

    _add_image_safe(s, "09_optical_composite_6k.jpg",
                    Inches(0.6), Inches(1.1), Inches(12.0), Inches(5.2))

    _add_multiline(s, Inches(0.6), Inches(6.5), Inches(12.0), Inches(0.8), [
        ("左：五谐振器 δf/f₀ vs 激光功率散点图。右：R4 (4.997 GHz) 在 −25/−30/−45 dBm 三功率下的 S21 局部放大 (±15 MHz)。三个功率下谐振谷形保持良好，谷底频率一致 → 读出功率不影响 f₀ 定位精度。",
         False, Pt(12), GRAY),
    ], line_spacing=Pt(4))
    _add_page_number(s, page[0], TOTAL)

    # ================================================================
    # Slide 12: 总结与展望
    # ================================================================
    s = new_slide()
    _set_slide_bg(s, 0xF8, 0xF8, 0xF8)

    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.08))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()

    _add_textbox(s, Inches(0.6), Inches(0.5), Inches(12.0), Inches(0.8),
                 "总结与展望", font_size=Pt(32), bold=True, color=DARK)

    # 三列布局
    col_width = Inches(3.8)
    col_gap = Inches(0.3)
    col_start = Inches(0.6)

    # --- 列 1: 已确认 ---
    done = [
        ("✓ 已确认", True, Pt(18), ACCENT),
        ("", False, Pt(6), DARK),
        ("5 枚谐振器 3.85–5.25 GHz", False, Pt(13), DARK),
        ("QL 2465–4167 @ 6 K", False, Pt(13), DARK),
        ("f₀(T) 一致红移 → 薄膜均匀", False, Pt(13), DARK),
        ("光学响应线性 → 准粒子拆对机制", False, Pt(13), DARK),
        ("QL vs P_read 无显著非线性", False, Pt(13), DARK),
        ("归一化 Δf/f₀(T) 高度重合", False, Pt(13), DARK),
    ]
    _add_multiline(s, col_start, Inches(1.8), col_width, Inches(4.5),
                   done, line_spacing=Pt(12))

    # --- 列 2: 关键发现 ---
    findings = [
        ("★ 关键发现", True, Pt(18), ACCENT),
        ("", False, Pt(6), DARK),
        ("R4 (4.997 GHz) 综合最优:", False, Pt(13), DARK),
        ("  最深耦合 −17.7 dB", False, Pt(13), DARK),
        ("  QL ≈ 3966", False, Pt(13), DARK),
        ("  可追踪至 ~74 K", False, Pt(13), DARK),
        ("  响应率最高 −1260 kHz @ 9 mW", False, Pt(13), DARK),
        ("", False, Pt(6), DARK),
        ("R3 (4.500 GHz) QL 最高 ≈ 4167", False, Pt(13), DARK),
        ("  Dip 较浅 (−6 dB) → 耦合偏弱", False, Pt(13), DARK),
        ("", False, Pt(6), DARK),
        ("R5 (5.252 GHz) 近带边 → 早消失", False, Pt(13), DARK),
    ]
    _add_multiline(s, col_start + col_width + col_gap, Inches(1.8), col_width, Inches(4.5),
                   findings, line_spacing=Pt(12))

    # --- 列 3: 下一步 ---
    next_steps = [
        ("→ 下一步", True, Pt(18), ACCENT),
        ("", False, Pt(6), DARK),
        ("Circle-fit 分离 Qc / Qi", False, Pt(13), DARK),
        ("变温光学响应 — 建立响应率 R(T)", False, Pt(13), DARK),
        ("NEP (噪声等效功率) 估算", False, Pt(13), DARK),
        ("多像素统计 — 器件均匀性", False, Pt(13), DARK),
        ("更低温度测量 (< 1 K)", False, Pt(13), DARK),
        ("与 BCS 理论定量比较", False, Pt(13), DARK),
        ("遴选最佳 KID 像素", False, Pt(13), DARK),
    ]
    _add_multiline(s, col_start + 2*(col_width + col_gap), Inches(1.8), col_width, Inches(4.5),
                   next_steps, line_spacing=Pt(12))

    _add_page_number(s, page[0], TOTAL)

    # ---- 保存 ----
    PPTX_OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(PPTX_OUT))
    return PPTX_OUT


# ============================================================
# 主入口
# ============================================================
if __name__ == "__main__":
    path = build_pptx()
    print(f"PPTX 已保存到: {path}")
    print(f"共 12 页幻灯片")
