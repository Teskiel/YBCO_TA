# -*- coding: utf-8 -*-
"""
Gurney 模板填充 — YBCO 五谐振器对比简报。
"""

import os, sys
from pathlib import Path

_script_dir = Path(__file__).resolve().parent
_otherwise_dir = _script_dir / "otherwise"
if str(_otherwise_dir) not in sys.path:
    sys.path.insert(0, str(_otherwise_dir))

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN

TEMPLATE = _script_dir / "templates" / "Gurney_clean.pptx"
FIG_DIR  = _script_dir / "output" / "merged" / "compare"
PPTX_OUT = _script_dir / "output" / "YBCO_KID_五谐振器对比_v5.pptx"

R = [
    {"id":"R1","f0":3.8460,"dip":-5.80,"ql":2465,"df9":-900},
    {"id":"R2","f0":4.0096,"dip":-4.33,"ql":2673,"df9":-960},
    {"id":"R3","f0":4.5002,"dip":-5.98,"ql":4167,"df9":-1140},
    {"id":"R4","f0":4.9970,"dip":-17.74,"ql":3966,"df9":-1260},
    {"id":"R5","f0":5.2516,"dip":-6.10,"ql":2501,"df9":-1260},
]

# 布局常量 (英寸) —— 全部手动控制, 不依赖模板占位符位置
TITLE_TOP   = 0.25   # 标题距顶
TITLE_H     = 0.45   # 标题高度
FIG_TOP     = 0.85   # 图片距顶
FIG_H       = 4.8    # 图片高度
TEXT_TOP    = 5.8    # 文字距顶
TEXT_H      = 1.5    # 文字高度
MARGIN      = 0.4    # 左右边距
FIG_W       = 9.2    # 图片宽度
TITLE_SIZE  = Pt(18) # 页标题字号
BODY_SIZE   = Pt(14) # 正文
SMALL_SIZE  = Pt(11) # 脚注


def fill_ph(slide, idx, text):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            ph.text = text
            return ph
    return None


def add_fig(slide, name, left, top, width, height=None):
    for ext in [".jpg", ".svg", ".png"]:
        p = FIG_DIR / (name + ext)
        if p.exists():
            return slide.shapes.add_picture(str(p), left, top, width, height)
    return None


def title_tb(slide, text):
    """在页面顶部加标题文本框。"""
    tb = slide.shapes.add_textbox(Inches(MARGIN), Inches(TITLE_TOP),
                                   Inches(FIG_W), Inches(TITLE_H))
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = TITLE_SIZE
    p.font.bold = True
    return tb


def body_tb(slide, paras):
    """在页面底部加正文文本框。paras: [(text, size, bold), ...]"""
    tb = slide.shapes.add_textbox(Inches(MARGIN), Inches(TEXT_TOP),
                                   Inches(FIG_W), Inches(TEXT_H))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (txt, sz, bld) in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt
        p.font.size = sz
        p.font.bold = bld
        p.space_after = Pt(6)
    return tb


def find_layout(prs, *names):
    for lo in prs.slide_layouts:
        if lo.name in names:
            return lo
    return prs.slide_layouts[6]


# ============================================================
def build():
    prs = Presentation(str(TEMPLATE))

    LAY_TITLE_SUB  = find_layout(prs, "TITLE_1")
    LAY_TITLE_BODY = find_layout(prs, "TITLE_AND_BODY")
    LAY_BLANK      = find_layout(prs, "BLANK")

    # ---- S1: 封面 (用模板自带 TITLE_1 布局) ----
    s = prs.slides.add_slide(LAY_TITLE_SUB)
    fill_ph(s, 0, "YBCO KID 片上五谐振器对比表征")
    fill_ph(s, 1, "赵思源\n频率跨度 3.85 – 5.25 GHz    6 K → 80 K    -25 dBm    2026-06-15")

    # ---- S2: 全谱指纹 (空白版式 + 手动标题 + 图 + 右侧表格) ----
    s = prs.slides.add_slide(LAY_BLANK)
    title_tb(s, "全谱指纹：五谐振器 S21 全景  @  6 K ,  -25 dBm ,  0 mW")
    add_fig(s, "01_spectrum_overview",
            Inches(MARGIN), Inches(FIG_TOP), Inches(7.8), Inches(FIG_H))

    # 右侧参数表
    tx = Inches(8.5)
    ty = Inches(FIG_TOP)
    tb = s.shapes.add_textbox(tx, ty, Inches(2.3), Inches(FIG_H))
    tf = tb.text_frame; tf.word_wrap = True
    rows = [
        ("参数  @ 6 K", Pt(12), True),
        ("", Pt(4), False),
        ("  f₀ / GHz     Dip / dB     QL", Pt(8), True),
        ("", Pt(2), False),
    ]
    for r in R:
        rows.append((f"{r['id']}   {r['f0']:.4f}   {r['dip']:+.1f}    {r['ql']}",
                     Pt(9), r['id'] == 'R4'))
    for _ in range(8):  # 空白行填充
        rows.append(("", Pt(7), False))

    for i, (txt, sz, bld) in enumerate(rows):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt
        p.font.size = sz
        p.font.bold = bld
        p.space_after = Pt(4)

    # ---- S3: f₀(T) (图在上, 文在下) ----
    s = prs.slides.add_slide(LAY_BLANK)
    title_tb(s, "谐振频率温度响应：f₀(T) 五线对比  @  -25 dBm ,  0 mW")
    add_fig(s, "02_f0_vs_temp_all",
            Inches(MARGIN), Inches(FIG_TOP), Inches(FIG_W), Inches(FIG_H))
    body_tb(s, [
        ("全部谐振器随温度单调红移, 频移幅度一致 → YBCO 薄膜均匀, Tc 一致。动能电感 Lₖ ∝ λ²(T) 驱动频移。",
         BODY_SIZE, False),
        (f"R1 – R4 可追踪至 ~70 K;  R5 (5.252 GHz) 近带边, 可追踪至 ~28 K。高温段 (>46 K) 谐振峰被热准粒子淹没。",
         SMALL_SIZE, False),
    ])

    # ---- S4: Qi(T) ----
    s = prs.slides.add_slide(LAY_BLANK)
    title_tb(s, "品质因子温度响应：QL(T) 五线对比  @  -25 dBm ,  0 mW")
    add_fig(s, "03_qi_vs_temp_all",
            Inches(MARGIN), Inches(FIG_TOP), Inches(FIG_W), Inches(FIG_H))
    body_tb(s, [
        (f"R3 最高 QL ≈ {R[2]['ql']} (@ 4.500 GHz);  R4 最深耦合 ({R[3]['dip']:.1f} dB), QL 仍保持 ≈ {R[3]['ql']}。",
         BODY_SIZE, False),
        ("QL 随温度升高下降 → 准粒子损耗增大。需 circle-fit 分离耦合品质因子 Qc 与内禀品质因子 Qi。",
         SMALL_SIZE, False),
    ])

    # ---- S5: 6K 光学响应 2×2 复合图 ----
    s = prs.slides.add_slide(LAY_BLANK)
    title_tb(s, "6 K 光学响应：δf / f₀  散点 + 三功率 S21 局部放大")
    add_fig(s, "09_optical_composite_6k",
            Inches(MARGIN), Inches(FIG_TOP), Inches(FIG_W), Inches(FIG_H))
    body_tb(s, [
        ("左：δf/f₀ vs 激光功率，纯散点。右：R4 (4.997 GHz) 在 -25/-30/-45 dBm 下的 S21 局部放大 (±15 MHz)，展示读出功率对谐振谷形的影响。",
         SMALL_SIZE, False),
    ])

    # ---- S6: 高温光学响应 2×2 复合图 ----
    s = prs.slides.add_slide(LAY_BLANK)
    title_tb(s, "高温光学响应：δf / f₀  随温度衰减 + 最高可测温度 S21 局部放大")
    add_fig(s, "10_optical_composite_hight",
            Inches(MARGIN), Inches(FIG_TOP), Inches(FIG_W), Inches(FIG_H))
    body_tb(s, [
        ("响应率随温度升高趋近于零 → 热准粒子淹没光注入信号。右：最高可测温度处 R4 的 S21 局部放大。",
         SMALL_SIZE, False),
    ])

    # ---- S7: QL vs 读出功率 ----
    s = prs.slides.add_slide(LAY_BLANK)
    title_tb(s, "QL  vs  VNA 读出功率  @  6 K ,  0 mW")
    add_fig(s, "05_ql_vs_power",
            Inches(MARGIN), Inches(FIG_TOP), Inches(FIG_W), Inches(FIG_H))
    body_tb(s, [
        ("三条 VNA 功率 (-45, -30, -25 dBm) 下各谐振器的 QL 对比。若 QL 随读出功率显著变化 → 非线性效应。",
         BODY_SIZE, False),
    ])

    # ---- S7: Δf/f₀ vs T（全温光学响应率）----
    s = prs.slides.add_slide(LAY_BLANK)
    title_tb(s, "归一化光学响应率：Δf / f₀  vs  温度  @  0 → 9 mW")
    add_fig(s, "06_responsivity_vs_temp",
            Inches(MARGIN), Inches(FIG_TOP), Inches(FIG_W), Inches(FIG_H))
    body_tb(s, [
        ("Δf/f₀ = [f₀(9 mW) - f₀(0 mW)] / f₀(0 mW)，单位 ppm。响应率随温度升高趋近于零 → 热准粒子淹没光注入信号。",
         BODY_SIZE, False),
        ("各谐振器曲线自然截断于各自超导态消失的温度：R4 可追踪至 ~74 K，R1–R3 至 ~40–58 K，R5 至 ~28 K。",
         SMALL_SIZE, False),
    ])

    # ---- S8: Dip 深度 vs T ----
    s = prs.slides.add_slide(LAY_BLANK)
    title_tb(s, "谐振谷深度  vs  温度")
    add_fig(s, "07_dip_vs_temp",
            Inches(MARGIN), Inches(FIG_TOP), Inches(FIG_W), Inches(FIG_H))
    body_tb(s, [
        ("虚线 = 检测极限 (dip = -1 dB)。谷深随 T 单调变浅 → 热准粒子增多 → 谐振展宽。",
         BODY_SIZE, False),
        ("R4 最深 (-17.7 dB @ 6K)，坚持最远 (~74K)；R5 最浅，最早消失 (~28K)。",
         SMALL_SIZE, False),
    ])

    # ---- S9: 归一化 f₀ 偏移 vs T ----
    s = prs.slides.add_slide(LAY_BLANK)
    title_tb(s, "归一化谐振频率偏移： (f₀(T) - f₀(6K)) / f₀(6K)  vs  温度")
    add_fig(s, "08_normalized_f0_vs_temp",
            Inches(MARGIN), Inches(FIG_TOP), Inches(FIG_W), Inches(FIG_H))
    body_tb(s, [
        ("五条归一化曲线在重叠温区高度重合 → YBCO 薄膜的动能电感温度响应高度一致，Tc 均匀。",
         BODY_SIZE, False),
        ("不同谐振器的绝对 f₀ 不同（几何差异），但相对频移 Δf/f₀(T) 由材料超导能隙 Δ(T) 决定，应一致。",
         SMALL_SIZE, False),
    ])

    # ---- S10: 总结 (用模板 TITLE_AND_BODY) ----
    s = prs.slides.add_slide(LAY_TITLE_BODY)
    fill_ph(s, 0, "总结与展望")
    fill_ph(s, 1,
        f"✓  5 枚谐振器覆盖 3.85 – 5.25 GHz, QL 范围 2465 – 4167 @ 6 K\n"
        f"✓  R3 (4.500 GHz) 品质因子最高 ≈ {R[2]['ql']};  R4 (4.997 GHz) 耦合最深 ({R[3]['dip']:.1f} dB)\n"
        f"✓  f₀(T) 一致红移 → YBCO 薄膜均匀, Tc 一致\n"
        f"✓  光学响应率与频率正相关: R4 最灵敏 ({R[3]['df9']:.2f} kHz @ 9 mW)\n"
        f"✓  QL 随 T 退化 → 热准粒子损耗主导高温行为\n"
        f"⚠  R5 (5.252 GHz) 近带边 → 追踪距离有限 (~28 K), 需优化设计\n\n"
        f"下一步: Circle-fit 分离 Qc/Qi → 变功率测量 → 多温度光学响应 → NEP 估算 → 遴选最佳像素"
    )

    return prs


if __name__ == "__main__":
    print("Building Gurney template PPTX...")
    prs = build()
    os.makedirs(PPTX_OUT.parent, exist_ok=True)
    prs.save(str(PPTX_OUT))
    print(f"Saved: {PPTX_OUT}")
