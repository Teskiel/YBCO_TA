# -*- coding: utf-8 -*-
"""测试 PPT 生成脚本的关键函数。"""
import sys
from pathlib import Path

_script_dir = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(_script_dir))

import os
import tempfile
from generate_ppt import (
    scan_temperatures,
    find_first_s2p,
    extract_key_numbers,
    build_pptx,
    MERGED_DIR,
)


def test_scan_temperatures():
    """应该返回排序后的温度列表。"""
    temps = scan_temperatures()
    assert len(temps) > 0, "应该至少找到一个温度目录"
    assert all(isinstance(t, tuple) and len(t) == 2 for t in temps), \
        "每个元素应为 (int_temp, dirname) 元组"
    # 验证升序
    int_temps = [t for t, _ in temps]
    assert int_temps == sorted(int_temps), "温度列表应为升序"
    assert int_temps[0] >= 4, f"最低温应 ≥ 4K: {int_temps[0]}"
    assert all(isinstance(t, int) for t in int_temps), "温度应为整数"


def test_find_first_s2p():
    """应该返回有效 S2P 文件路径。"""
    temp_entries = scan_temperatures()
    _, dirname = temp_entries[0]
    path = find_first_s2p(dirname, 25, 0)
    assert path is not None, f"应该找到 {dirname}/-25dBm/00mW 的 S2P 文件"
    assert os.path.isfile(path), f"应为有效文件: {path}"
    assert path.endswith(".s2p"), f"应以 .s2p 结尾: {path}"


def test_extract_key_numbers():
    """应该返回所有关键字段的非空值。"""
    nums = extract_key_numbers()

    required_keys = [
        "f0_ghz", "num_resonances", "low_temp", "high_temp",
        "all_resonances",
    ]
    for key in required_keys:
        assert key in nums, f"缺少字段: {key}"
        assert nums[key] is not None, f"字段 {key} 为 None"

    assert nums["num_resonances"] > 0, "应该检测到至少 1 个谐振峰"
    assert nums["f0_ghz"] > 0, "谐振频率应为正数"
    assert nums["low_temp"] < nums["high_temp"], "低温应小于高温"


def test_build_pptx_creates_file():
    """build_pptx 应该生成有效的 PPTX 文件。"""
    nums = extract_key_numbers()

    with tempfile.TemporaryDirectory() as tmpdir:
        out_path = Path(tmpdir) / "test_output.pptx"
        prs = build_pptx(nums)
        prs.save(str(out_path))

        assert out_path.exists(), "PPTX 文件应该存在"
        assert out_path.stat().st_size > 1000, "文件不应为空"
        assert out_path.suffix == ".pptx"

        # 验证幻灯片数量
        from pptx import Presentation
        prs_check = Presentation(str(out_path))
        assert len(prs_check.slides) == 7, \
            f"应有 7 页幻灯片: 实际 {len(prs_check.slides)}"
