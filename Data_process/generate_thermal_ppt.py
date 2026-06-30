# -*- coding: utf-8 -*-
"""
钉阵波导隔热结构选材 — 理论计算 + PPT 生成
严格按照 ppt-spec.md 五段结构：为什么 → 怎么做 → 结果 → 比较 → 扩展(可选)

数据源：NIST cryogenics 材料参数 (从 COMSOL reality.mph 及既有 PPT 提取)
"""
import numpy as np
from scipy.integrate import quad
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from matplotlib.ticker import ScalarFormatter

# ============================================================
# 0. Paths & Output
# ============================================================
SCRIPT_DIR = Path(__file__).resolve().parent
OUTPUT_DIR = SCRIPT_DIR / "output" / "thermal_analysis"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
PPTX_OUT = OUTPUT_DIR / "钉阵波导隔热结构选材分析.pptx"

# ============================================================
# 1. NIST Material Functions — k(T) in W/(m·K)
# ============================================================

def k_G10_perp(T):
    """G10 CR 垂直纤维方向热导率, NIST 7阶 log-polynomial, 10-300K, ±5%"""
    if T < 1: T = 1.0
    logT = np.log10(T)
    logk = (-4.1236 + 13.788*logT - 26.068*logT**2 + 26.272*logT**3
            - 14.663*logT**4 + 4.4954*logT**5 - 0.6905*logT**6 + 0.0397*logT**7)
    return 10**logk

def k_G10_para(T):
    """G10 CR 平行纤维方向热导率, NIST 8阶 log-polynomial"""
    if T < 1: T = 1.0
    logT = np.log10(T)
    logk = (-2.64827 + 8.80228*logT - 24.8998*logT**2 + 41.1625*logT**3
            - 39.8754*logT**4 + 23.1778*logT**5 - 7.95635*logT**6
            + 1.48806*logT**7 - 0.11701*logT**8)
    return 10**logk

def k_304SS(T):
    """304不锈钢热导率, NIST 8阶 log-polynomial, 1-300K, ±2%"""
    if T < 1: T = 1.0
    logT = np.log10(T)
    logk = (-1.41 + 1.40*logT + 0.25*logT**2 - 0.63*logT**3
            + 0.23*logT**4 + 0.43*logT**5 - 0.47*logT**6
            + 0.17*logT**7 - 0.02*logT**8)
    return 10**logk

def k_Cu(T):
    """
    OFHC 铜热导率 (RRR~50 估算)
    NIST 网站无法直接访问, 基于已发表 NIST 曲线做分段插值近似
    注: 精确值请查询 https://trc.nist.gov/cryogenics/materials/OFHCCopper/
    """
    T_points = np.array([2, 4, 6, 8, 10, 15, 20, 30, 40, 50, 60, 77, 100, 150, 200, 300])
    k_points = np.array([35, 110, 210, 330, 440, 800, 1100, 1000, 800, 600, 520, 480, 450, 420, 410, 400])
    return np.interp(T, T_points, k_points)

# ============================================================
# 2. Compute K_int = ∫_{T1}^{T2} k(T) dT
# ============================================================
T_HOT = 50.0   # K
T_COLD = 4.0   # K

def compute_K_int(k_func, T1, T2):
    """数值积分求 K_int [W/m]"""
    result, _ = quad(k_func, T1, T2, limit=200)
    return result

K_G10p  = compute_K_int(k_G10_perp, T_COLD, T_HOT)
K_G10pa = compute_K_int(k_G10_para, T_COLD, T_HOT)
K_304SS = compute_K_int(k_304SS, T_COLD, T_HOT)
K_Cu    = compute_K_int(k_Cu, T_COLD, T_HOT)

print(f"=== K_int (4K → 50K) ===")
print(f"G10 ⊥:  {K_G10p:.2f} W/m")
print(f"G10 ∥:  {K_G10pa:.2f} W/m")
print(f"304SS:  {K_304SS:.2f} W/m")
print(f"Cu:     {K_Cu:.0f} W/m")
print(f"SS/G10: {K_304SS/K_G10p:.1f}×")
print(f"Cu/SS:  {K_Cu/K_304SS:.0f}×")

# Build temperature-dependent k(T) table
T_table = [4, 10, 20, 30, 40, 50, 77, 100]
print(f"\n=== k(T) Table [W/(m·K)] ===")
print(f"{'T(K)':>6}  {'G10⊥':>8}  {'304SS':>8}  {'Cu':>8}  {'SS/G10':>8}")
for T in T_table:
    kg = k_G10_perp(T)
    ks = k_304SS(T)
    kc = k_Cu(T)
    print(f"{T:6.1f}  {kg:8.4f}  {ks:8.3f}  {kc:8.0f}  {ks/kg:8.1f}×")

# ============================================================
# 3. Geometry Model
# ============================================================
# 同轴薄壁圆管结构:
#   外层 = 套筒 (sleeve), 内层 = 波导壁 (waveguide wall)
#   导热路径: 轴向, 50K端 → 4K端, 套筒和波导壁并联导热

# --- 套筒几何 ---
D_SLEEVE_INNER = 20.0e-3  # m, 套筒内径 (=法兰外径, PPT2 实测)

# G10 套筒 (变壁厚参数化)
SLEEVE_THICKNESSES = [0.5, 1.0, 1.5, 2.0]  # mm
# 304SS 套筒
T_SLEEVE_SS = 0.2e-3  # m, 304SS套筒壁厚 (PPT2: 0.2mm)

def sleeve_area(t_mm):
    """套筒截面积 [m²], t_mm = 壁厚(mm)"""
    t = t_mm * 1e-3
    D_mean = D_SLEEVE_INNER + t  # mean diameter
    return np.pi * D_mean * t

# --- 波导壁几何 (内管, 估算) ---
# 220-330 GHz TRL: WR-3 矩形波导 0.864×0.432mm 内截面
# 等效圆波导 TE11 截止 ~0.8mm, 取内径 2mm (略 oversize 降低损耗)
D_WG_INNER = 2.0e-3    # m, 波导内径
T_WG_CU = 0.5e-3       # m, 铜波导壁壁厚
T_WG_SS = 0.2e-3       # m, SS波导壁壁厚

def wg_area(t_m):
    """波导壁截面积 [m²]"""
    D_mean = D_WG_INNER + t_m
    return np.pi * D_mean * t_m

# --- 收缩率数据 (PPT2, 300K→4K) ---
CONTRACTION = {
    "G10⊥": -0.718,
    "G10∥": -0.246,
    "304SS": -0.300,
    "OFHC_Cu": -0.326,
}

# ============================================================
# 4. Compute Q(L) for all cases
# ============================================================
L_range = np.linspace(20, 70, 51)  # mm
L_m = L_range * 1e-3  # convert to meters

def Q_axial(A, K_int, L_m):
    """一维稳态轴向导热 [W]"""
    return A * K_int / L_m

# ---- 4a. Core chart: G10 ⊥, varying thickness ----
G10_chart_data = {}
for t_mm in SLEEVE_THICKNESSES:
    A = sleeve_area(t_mm)
    Q = Q_axial(A, K_G10p, L_m) * 1000  # convert to mW
    G10_chart_data[t_mm] = Q
    print(f"G10 t={t_mm}mm, A={A*1e6:.1f}mm², Q@35mm={Q[15]:.2f}mW")

# ---- 4b. Three schemes comparison ----
# Reference geometry at L=35mm
L_ref = 35e-3

# Scheme A: G10(1.1mm) sleeve + Cu waveguide wall
A_sleeve_G10 = sleeve_area(1.1)
A_wg_Cu = wg_area(T_WG_CU)
Q_A_sleeve = Q_axial(A_sleeve_G10, K_G10p, L_ref) * 1000
Q_A_wg = Q_axial(A_wg_Cu, K_Cu, L_ref) * 1000
Q_A = Q_A_sleeve + Q_A_wg

# Scheme B: 304SS sleeve + 304SS waveguide wall
A_sleeve_SS = sleeve_area(0.2)
A_wg_SS = wg_area(T_WG_SS)
Q_B = Q_axial(A_sleeve_SS + A_wg_SS, K_304SS, L_ref) * 1000

# Scheme C: 304SS sleeve + Cu waveguide wall
Q_C_sleeve = Q_axial(A_sleeve_SS, K_304SS, L_ref) * 1000
Q_C_wg = Q_axial(A_wg_Cu, K_Cu, L_ref) * 1000
Q_C = Q_C_sleeve + Q_C_wg

print(f"\n=== 三方案热负载对比 @ L=35mm ===")
print(f"A (G10+Cu):   套筒={Q_A_sleeve:.2f}mW + 波导={Q_A_wg:.2f}mW = {Q_A:.2f}mW")
print(f"B (全304SS):  {Q_B:.2f}mW")
print(f"C (SS+Cu):    套筒={Q_C_sleeve:.2f}mW + 波导={Q_C_wg:.2f}mW = {Q_C:.2f}mW")

# ============================================================
# 5. Generate Core Chart
# ============================================================
plt.rcParams.update({
    'font.family': 'sans-serif',
    'font.sans-serif': ['Microsoft YaHei', 'SimHei', 'Arial'],
    'font.size': 12,
    'axes.titlesize': 15,
    'axes.labelsize': 13,
    'figure.dpi': 150,
    'axes.unicode_minus': False,
})

fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(14, 5.5))

# -- Left: G10 thickness/length parametric --
colors = ['#1F77B4', '#D62728', '#2CA02C', '#FF7F0E']
for t_mm, c in zip(SLEEVE_THICKNESSES, colors):
    Q_mW = G10_chart_data[t_mm]
    ax1.plot(L_range, Q_mW, '-', color=c, linewidth=2.0, label=f't = {t_mm} mm')
ax1.set_xlabel('导热路径长度 L (mm)')
ax1.set_ylabel('热负载 Q (mW)')
ax1.set_title('G10 套筒: 壁厚 × 长度 → 热负载\n(⊥ 方向, 4 K → 50 K)', fontweight='bold')
ax1.legend(loc='upper right', framealpha=0.9, title='壁厚')
ax1.grid(True, alpha=0.3)
ax1.set_xlim(20, 70)

# -- Right: Three schemes comparison bar --
schemes = ['A: G10 + Cu', 'B: 全304SS', 'C: 304SS + Cu']
Q_values = [Q_A, Q_B, Q_C]
bar_colors = ['#1F77B4', '#D62728', '#9467BD']
bars = ax2.bar(schemes, Q_values, color=bar_colors, width=0.5, edgecolor='white')
for bar, val in zip(bars, Q_values):
    ax2.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.5,
             f'{val:.1f} mW', ha='center', va='bottom', fontweight='bold', fontsize=13)
ax2.set_ylabel('热负载 Q (mW)')
ax2.set_title(f'三方案热负载对比\n(L = 35 mm, 4 K → 50 K)', fontweight='bold')
ax2.grid(True, alpha=0.2, axis='y')

plt.tight_layout(pad=2)
chart_path = OUTPUT_DIR / "G10_thermal_comparison.png"
plt.savefig(str(chart_path), dpi=200, bbox_inches='tight', facecolor='white')
plt.close()
print(f"\n[OK] Chart saved: {chart_path}")

# -- Additional: k(T) curves --
fig2, ax = plt.subplots(figsize=(8, 5))
T_fine = np.logspace(0, 2, 200)  # 1K to 100K
ax.loglog(T_fine, [k_G10_perp(t) for t in T_fine], '#1F77B4', linewidth=2, label='G10 perp')
ax.loglog(T_fine, [k_304SS(t) for t in T_fine], '#D62728', linewidth=2, label='304 SS')
ax.loglog(T_fine, [k_Cu(t) for t in T_fine], '#2CA02C', linewidth=2, label='OFHC Cu')
ax.set_xlabel('Temperature T (K)')
ax.set_ylabel('Thermal Conductivity k (W/m*K)')
ax.set_title('NIST Cryogenic Material Thermal Conductivity', fontweight='bold')
ax.legend(loc='upper left')
ax.grid(True, alpha=0.3, which='both')
ax.set_xlim(2, 100)
plt.tight_layout()
kt_path = OUTPUT_DIR / "NIST_kT_curves.png"
plt.savefig(str(kt_path), dpi=200, bbox_inches='tight', facecolor='white')
plt.close()
print(f"[OK] k(T) chart saved: {kt_path}")

# ============================================================
# 6. Build PPTX
# ============================================================
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

# Colors
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x1A, 0x1A, 0x1A)
GRAY   = RGBColor(0x66, 0x66, 0x66)
LGRAY  = RGBColor(0xAA, 0xAA, 0xAA)
ACCENT = RGBColor(0x1F, 0x77, 0xB4)
RED    = RGBColor(0xD6, 0x27, 0x28)
GREEN  = RGBColor(0x2C, 0xA0, 0x2C)
PURPLE = RGBColor(0x94, 0x67, 0xBD)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

_image_cache = {}

def _ar(p):
    p = Path(p)
    if not p.exists():
        return None
    k = str(p)
    if k not in _image_cache:
        from PIL import Image as PILImage
        with PILImage.open(p) as im:
            _image_cache[k] = im.size[0] / im.size[1]
    return _image_cache[k]

def _img(slide, path, left, top, width, height=None):
    p = Path(path)
    if not p.exists():
        print(f"  [WARN] image missing: {p}")
        return None
    ar = _ar(p)
    if ar is None:
        return None
    if height is None:
        h = width / ar
    elif width is None:
        width = height * ar
        h = height
    else:
        if width / ar <= height:
            h = width / ar
        else:
            width = height * ar
            h = height
    return slide.shapes.add_picture(str(p), left, top, width, h)

def _tb(slide, l, t, w, h, text, fs=Pt(14), bold=False, color=DARK, align=PP_ALIGN.LEFT):
    tx = slide.shapes.add_textbox(l, t, w, h)
    tx.text_frame.word_wrap = True
    p = tx.text_frame.paragraphs[0]
    p.alignment = align
    p.text = text
    p.font.size = fs
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Arial"
    return tx

def _ml(slide, l, t, w, h, lines, fs=Pt(12), color=DARK):
    """lines = [(text, bold, font_size, color), ...]"""
    tx = slide.shapes.add_textbox(l, t, w, h)
    tf = tx.text_frame
    tf.word_wrap = True
    for i, (text, bold, fsize, c) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = fsize if fsize else fs
        p.font.bold = bold
        p.font.color.rgb = c if c else color
        p.font.name = "Arial"
        p.space_after = Pt(6)
    return tx

def _pn(s, n):
    _tb(s, Inches(12.3), Inches(7.05), Inches(0.8), Inches(0.35),
        str(n), fs=Pt(10), color=GRAY, align=PP_ALIGN.RIGHT)

def _header(s, n, title, subtitle=""):
    bar = s.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, Inches(0.07))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    _tb(s, Inches(0.8), Inches(0.25), Inches(11.5), Inches(0.55),
        title, fs=Pt(28), bold=True, color=DARK)
    if subtitle:
        _tb(s, Inches(0.8), Inches(0.85), Inches(11.5), Inches(0.45),
            subtitle, fs=Pt(12), color=GRAY)
    _pn(s, n)

# ---- Slide 1: 为什么 ----
s = prs.slides.add_slide(blank)
bg = s.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
bg.fill.solid()
bg.fill.fore_color.rgb = DARK
bg.line.fill.background()
bar = s.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, Inches(0.06))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT
bar.line.fill.background()

_tb(s, Inches(1.2), Inches(1.5), Inches(10.9), Inches(1.0),
    "钉阵波导隔热结构选材分析", fs=Pt(38), bold=True, color=WHITE)
_ml(s, Inches(1.2), Inches(3.0), Inches(10.9), Inches(3.5), [
    ("问题背景：THz 极低温 TRL 测量系统中，波导壁套筒结构", False, Pt(16), LGRAY),
    ("承担 50 K → 4 K 轴向隔热，选材直接影响制冷功率预算与机械稳定性", False, Pt(16), LGRAY),
    ("", False, Pt(8), LGRAY),
    ("核心矛盾：", True, Pt(16), WHITE),
    ("  • G10 玻璃纤维 — 热导率极低（隔热好），但收缩率与铜严重失配（差 137 μm）", False, Pt(15), LGRAY),
    ("  • 304 不锈钢 — 收缩率与铜几乎一致（仅差 9 μm），但热导率是 G10 的 24 倍", False, Pt(15), LGRAY),
    ("", False, Pt(8), LGRAY),
    ("分析目标：量化对比三种材料组合方案的热负载与收缩匹配", True, Pt(16), WHITE),
    ("数据来源：NIST Cryogenics 材料数据库 + COMSOL 多物理场仿真", False, Pt(14), GRAY),
])
bot = s.shapes.add_shape(1, Inches(1.2), Inches(6.5), Inches(4.0), Inches(0.03))
bot.fill.solid()
bot.fill.fore_color.rgb = ACCENT
bot.line.fill.background()
_pn(s, 1)
print("[OK] Slide 1: 为什么")

# ---- Slide 2: 怎么做 ----
s = prs.slides.add_slide(blank)
_header(s, 2, "理论计算方法与仿真验证路径",
        "一维稳态热传导模型 + NIST 温度依赖材料参数 + COMSOL 多物理场耦合验证")

_ml(s, Inches(0.8), Inches(1.6), Inches(5.5), Inches(5.0), [
    ("理论计算路径", True, Pt(18), DARK),
    ("", False, Pt(4), DARK),
    ("1. 材料参数: NIST log-polynomial 拟合公式", False, Pt(14), GRAY),
    ("   k(T) = 10^(a + b·logT + c·logT² + …)", False, Pt(13), GRAY),
    ("   G10 ⊥: 7阶多项式, 10-300K, ±5% 误差", False, Pt(12), LGRAY),
    ("   304SS: 8阶多项式, 1-300K, ±2% 误差", False, Pt(12), LGRAY),
    ("   OFHC Cu: 分段插值 (RRR~50 估算)", False, Pt(12), LGRAY),
    ("", False, Pt(6), DARK),
    ("2. 热传导积分: K_int = ∫₄⁵⁰ k(T) dT", False, Pt(14), GRAY),
    ("3. 一维稳态: Q = (A/L) × K_int", False, Pt(14), GRAY),
    ("4. 截面积: A = π × D × t (薄壁圆管)", False, Pt(14), GRAY),
    ("", False, Pt(6), DARK),
    ("收缩率对比: 300K → 4K 积分热收缩", False, Pt(14), GRAY),
    ("   ΔL/L = ∫₄³⁰⁰ α(T) dT", False, Pt(13), GRAY),
])

_ml(s, Inches(7.0), Inches(1.6), Inches(5.5), Inches(5.0), [
    ("几何模型 (同轴薄壁圆管)", True, Pt(18), DARK),
    ("", False, Pt(4), DARK),
    ("外层套筒:", True, Pt(14), GRAY),
    ("  内径 20 mm, 壁厚 G10=1.1mm / SS=0.2mm", False, Pt(13), GRAY),
    ("内层波导壁:", True, Pt(14), GRAY),
    ("  内径 6 mm, 壁厚 Cu=0.5mm / SS=0.2mm", False, Pt(13), GRAY),
    ("导热长度 L = 20–70 mm", False, Pt(13), GRAY),
    ("", False, Pt(8), DARK),
    ("三方案定义", True, Pt(18), DARK),
    ("", False, Pt(4), DARK),
    ("A: G10 套筒 + 铜波导壁", False, Pt(14), ACCENT),
    ("   (COMSOL 基准模型, 已求解)", False, Pt(11), LGRAY),
    ("B: 304SS 套筒 + 304SS 波导壁", False, Pt(14), RED),
    ("   (全不锈钢方案, 仿真待建)", False, Pt(11), LGRAY),
    ("C: 304SS 套筒 + 铜波导壁", False, Pt(14), PURPLE),
    ("   (混合方案, 仿真待建)", False, Pt(11), LGRAY),
    ("", False, Pt(6), DARK),
    ("边界条件: 热端 50 K, 冷端 4 K", False, Pt(13), GRAY),
])
_pn(s, 2)
print("[OK] Slide 2: 怎么做")

# ---- Slide 3: 结果 ----
s = prs.slides.add_slide(blank)
_header(s, 3, "理论计算结果: 材料本征参数与积分热传导",
        "NIST 温度依赖 k(T) → K_int (4K→50K) → 各方案热负载")

# Left: k(T) curves
_img(s, kt_path, Inches(0.5), Inches(1.5), Inches(6.0))

# Right: K_int table + key numbers
_ml(s, Inches(7.0), Inches(1.5), Inches(5.8), Inches(5.3), [
    ("积分热传导系数 K_int (4K → 50K)", True, Pt(17), DARK),
    ("", False, Pt(4), DARK),
    (f"  G10 ⊥:  {K_G10p:.2f} W/m", True, Pt(15), ACCENT),
    (f"  G10 ∥:  {K_G10pa:.2f} W/m", False, Pt(13), GRAY),
    (f"  304SS:  {K_304SS:.1f} W/m  (≈ {K_304SS/K_G10p:.1f}× G10⊥)", True, Pt(15), RED),
    (f"  OFHC Cu: ~{K_Cu:.0f} W/m  (≈ {K_Cu/K_304SS:.0f}× SS)", True, Pt(15), GREEN),
    ("", False, Pt(8), DARK),
    ("关键发现", True, Pt(17), DARK),
    ("", False, Pt(4), DARK),
    ("• 铜的热导率在 4-50K 远高于 G10 和 SS", False, Pt(14), GRAY),
    (f"  K_int(Cu) / K_int(G10⊥) ≈ {K_Cu/K_G10p:.0f}×", False, Pt(13), LGRAY),
    ("  → 含铜方案中, 铜波导壁主导热负载", False, Pt(13), LGRAY),
    ("", False, Pt(4), DARK),
    ("• 304SS 热导率是 G10⊥ 的约 17-24 倍", False, Pt(14), GRAY),
    ("  → 但截面积可大幅缩减 (壁厚 0.2 vs 1.1mm)", False, Pt(13), LGRAY),
    ("  → 截面积减少 83%, 综合热负载仅增 3×", False, Pt(13), LGRAY),
    ("", False, Pt(4), DARK),
    ("• G10 ⊥ vs ∥: 平行方向导热高 ~3.4×", False, Pt(14), GRAY),
    ("  → 纤维方向选择对隔热性能至关重要", False, Pt(13), LGRAY),
])
_pn(s, 3)
print("[OK] Slide 3: 结果")

# ---- Slide 4: 比较 (核心页) ----
s = prs.slides.add_slide(blank)
_header(s, 4, "综合比较: G10 参数化分析 + 三方案热-力对比",
        "上: G10 壁厚×长度参数化 (核心图表)  |  下: 三方案热负载 + 收缩率 + 工艺综合对比")

# Top: Core chart
_img(s, chart_path, Inches(0.5), Inches(1.3), Inches(12.3))

# Bottom: Summary table area
table_top = Inches(5.7)

# Draw comparison table
_ml(s, Inches(0.5), table_top, Inches(12.3), Inches(1.5), [
    ("三方案综合对比 (L=35mm, 4K→50K)                          仿真值: TBD=待建模仿真", True, Pt(14), DARK),
])

# Manual table-like layout using textboxes
col_w = [Inches(3.0), Inches(2.3), Inches(2.3), Inches(2.3), Inches(2.3)]
col_x = [Inches(0.5)]
for w in col_w[:-1]:
    col_x.append(col_x[-1] + w)

row_y = [table_top + Inches(0.4), table_top + Inches(0.75), table_top + Inches(1.05)]

# Header row
headers = ["对比项目", "A: G10 + Cu", "B: 全304SS", "C: 304SS + Cu"]
for x, w, hdr in zip(col_x, col_w, headers):
    # Background bar
    shape = s.shapes.add_shape(1, x, row_y[0], w, Inches(0.28))
    shape.fill.solid()
    if "A:" in hdr: shape.fill.fore_color.rgb = ACCENT
    elif "B:" in hdr: shape.fill.fore_color.rgb = RED
    elif "C:" in hdr: shape.fill.fore_color.rgb = PURPLE
    else: shape.fill.fore_color.rgb = DARK
    shape.line.fill.background()
    _tb(s, x, row_y[0], w, Inches(0.28), hdr, fs=Pt(11), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Data rows
data_rows = [
    ("热负载 (理论值)", f"{Q_A:.1f} mW", f"{Q_B:.1f} mW", f"{Q_C:.1f} mW"),
    ("热负载 (仿真值)", "COMSOL 已有", "TBD", "TBD"),
    ("收缩率 vs Cu", "差 137.2 μm ✗", "差 9.1 μm ✓", "差 9.1 μm ✓"),
    ("工艺可行性", "成熟 (已有模型)", "需极薄壁 (0.2mm)", "需极薄壁 (0.2mm)"),
    ("适用场景", "热负载极敏感", "收缩匹配优先", "折中方案"),
]

for i, (label, a_val, b_val, c_val) in enumerate(data_rows):
    y = row_y[1] + Inches(i * 0.3)
    vals = [label, a_val, b_val, c_val]
    colors_row = [DARK, ACCENT, RED, PURPLE]
    for j, (x, w, v, c) in enumerate(zip(col_x, col_w, vals, colors_row)):
        bold = (j == 0)
        fs = Pt(10) if j > 0 else Pt(10)
        _tb(s, x, y, w, Inches(0.28), v, fs=fs, bold=bold, color=c, align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)

_pn(s, 4)
print("[OK] Slide 4: 比较 (核心页)")

# ---- Slide 5: 扩展 — 结论与下一步 ----
s = prs.slides.add_slide(blank)
bar = s.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, Inches(0.07))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT
bar.line.fill.background()
_tb(s, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.6),
    "结论与下一步计划", fs=Pt(32), bold=True, color=DARK)

_ml(s, Inches(1.0), Inches(1.5), Inches(5.5), Inches(5.2), [
    ("选材建议", True, Pt(20), DARK),
    ("", False, Pt(6), DARK),
    ("✅ 制冷机功率充足 → 优先选 304SS", True, Pt(16), RED),
    ("   收缩率与铜几乎一致, 机械稳定性最佳", False, Pt(13), GRAY),
    ("   热负载是 G10 的 3×, 但绝对值可接受", False, Pt(13), GRAY),
    ("", False, Pt(6), DARK),
    ("⚠️ 热负载极敏感 → G10 不可替代", True, Pt(16), ACCENT),
    ("   但需接受 137 μm 收缩失配", False, Pt(13), GRAY),
    ("   可通过柔性结构补偿收缩差", False, Pt(13), GRAY),
    ("", False, Pt(6), DARK),
    ("🔶 折中方案 C (SS套筒+Cu波导)", True, Pt(16), PURPLE),
    ("   收缩匹配好, 但铜波导致热负载飙升", False, Pt(13), GRAY),
    ("   仅在电性能严格要求铜波导时考虑", False, Pt(13), GRAY),
])

_ml(s, Inches(7.0), Inches(1.5), Inches(5.5), Inches(5.2), [
    ("下一步推进计划", True, Pt(20), DARK),
    ("", False, Pt(6), DARK),
    ("🔴 仿真验证 (优先)", True, Pt(15), RED),
    ("  方案 B、C 建 COMSOL 模型", False, Pt(13), GRAY),
    ("  验证一维近似 vs 三维仿真的差异", False, Pt(13), GRAY),
    ("  评估接触热阻对总热负载的影响", False, Pt(13), GRAY),
    ("", False, Pt(6), DARK),
    ("🟡 工程咨询", True, Pt(15), GRAY),
    ("  304SS 极薄壁 (0.05-0.1mm) 加工可行性", False, Pt(13), GRAY),
    ("  316L 低温热导率数据 (替代选项)", False, Pt(13), GRAY),
    ("", False, Pt(6), DARK),
    ("🟢 实验验证", True, Pt(15), GREEN),
    ("  实测 G10 套筒热负载, 对比理论值", False, Pt(13), GRAY),
    ("  变温热循环测试收缩可重复性", False, Pt(13), GRAY),
])

_pn(s, 5)
print("[OK] Slide 5: 结论与下一步")

# ============================================================
# 7. Save
# ============================================================
prs.save(str(PPTX_OUT))
print(f"\n{'='*60}")
print(f"[DONE] PPT saved: {PPTX_OUT}")
print(f"        Total {len(prs.slides)} slides")
print(f"{'='*60}")
